#!/usr/bin/env python3
"""Republican-newspaper PPTX slide deck generator."""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slides_spec import DECK, DISPLAY_TOTAL


PAPER = RGBColor(0xE8, 0xDF, 0xC9)
PAPER_LIGHT = RGBColor(0xF8, 0xF1, 0xDD)
PAPER_DEEP = RGBColor(0xDE, 0xD4, 0xBD)
PANEL = RGBColor(0xF5, 0xEB, 0xD3)
INK = RGBColor(0x16, 0x14, 0x11)
NEAR_BLACK = RGBColor(0x1F, 0x1B, 0x16)
DARK_WARM = RGBColor(0x3E, 0x38, 0x2E)
CHARCOAL = RGBColor(0x51, 0x4A, 0x3F)
OLIVE = RGBColor(0x6A, 0x62, 0x57)
STONE = RGBColor(0x8A, 0x81, 0x73)
RULE = RGBColor(0xB9, 0xAA, 0x91)
RULE_INK = RGBColor(0x4B, 0x40, 0x36)
IVORY = RGBColor(0xF8, 0xF1, 0xDD)
RED = RGBColor(0xB7, 0x35, 0x2A)

FONT_CONFIG = Path(__file__).resolve().parent.parent / "fonts" / "fonts.json"


def load_font_config() -> dict[str, dict[str, str]]:
    return json.loads(FONT_CONFIG.read_text(encoding="utf-8"))


def font_slot(role: str, slot: str = "latin") -> str:
    return load_font_config()[role][slot]


SERIF = font_slot("serif", "latin")
SERIF_EA = font_slot("serif", "eastAsian")
SANS = font_slot("sans", "latin")
MONO = font_slot("mono", "latin")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SHEET_X = Inches(0.34)
SHEET_Y = Inches(0.26)
SHEET_W = SLIDE_W - SHEET_X * 2
SHEET_H = SLIDE_H - SHEET_Y * 2
TEXTURE = Path(__file__).resolve().parent.parent / "images" / "paper-overlay.png"


def apply_typeface(run, font):
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    for tag in (
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
    ):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", SERIF_EA)
    cs = OxmlElement("a:cs")
    cs.set("typeface", font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def add_shape(slide, shape_type, left, top, width, height, fill=None, line=None, weight=1.0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(weight)
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill=None, line=None, weight=1.0):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill, line, weight)


def add_oval(slide, left, top, size, fill=None, line=None, weight=1.0):
    return add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill, line, weight)


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font=SANS,
    size=14,
    color=NEAR_BLACK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.0,
    tracking=False,
    rotation: float | None = None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    if rotation is not None:
        box.rotation = rotation
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text.upper() if tracking else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    apply_typeface(run, font)
    return box


def add_line(slide, x1, y1, x2, y2, color=INK, weight=1.0):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(weight)
    connector.shadow.inherit = False
    return connector


def verticalize(text: str) -> str:
    parts = []
    for chunk in text.split("\n"):
        parts.append("\n".join(char for char in chunk))
    return "\n\n".join(parts)


def patch_theme_fonts(pptx_path: str):
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = text.replace('<a:ea typeface=""/>', f'<a:ea typeface="{SERIF_EA}"/>')
                text = text.replace('script="Hans" typeface="宋体"', f'script="Hans" typeface="{SERIF_EA}"')
                text = text.replace('script="Hant" typeface="新細明體"', f'script="Hant" typeface="{SERIF_EA}"')
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def page_base(prs: Presentation, section: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    if TEXTURE.exists():
        slide.shapes.add_picture(str(TEXTURE), 0, 0, width=SLIDE_W, height=SLIDE_H)
    add_rect(slide, SHEET_X, SHEET_Y, SHEET_W, SHEET_H, None, INK, 1.15)
    add_rect(slide, SHEET_X + Inches(0.12), SHEET_Y + Inches(0.12), SHEET_W - Inches(0.24), SHEET_H - Inches(0.24), None, RULE_INK, 0.5)
    add_rect(slide, SHEET_X + Inches(0.22), SHEET_Y + Inches(0.22), SHEET_W - Inches(0.44), SHEET_H - Inches(0.44), None, RULE, 0.4)
    add_text(
        slide,
        section,
        SHEET_X + Inches(0.22),
        SHEET_Y + SHEET_H - Inches(0.34),
        Inches(4.2),
        Inches(0.16),
        font=MONO,
        size=7.2,
        color=STONE,
        tracking=True,
    )
    add_text(
        slide,
        f"{page:02d} / {DISPLAY_TOTAL:02d}",
        SHEET_X + SHEET_W - Inches(1.4),
        SHEET_Y + SHEET_H - Inches(0.34),
        Inches(1.05),
        Inches(0.16),
        font=MONO,
        size=7.2,
        color=STONE,
        align=PP_ALIGN.RIGHT,
        tracking=True,
    )
    return slide


def draw_issue_mark(slide, left, top, size, title, note):
    add_oval(slide, left, top, size, None, RED, 1.0)
    add_oval(slide, left + Inches(0.09), top + Inches(0.09), size - Inches(0.18), None, RED, 0.45)
    add_text(
        slide,
        title,
        left + Inches(0.08),
        top + Inches(0.25),
        size - Inches(0.16),
        Inches(0.34),
        font=SERIF,
        size=16,
        color=RED,
        bold=False,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        note,
        left + Inches(0.1),
        top + size - Inches(0.34),
        size - Inches(0.2),
        Inches(0.16),
        font=SANS,
        size=5.8,
        color=RED,
        align=PP_ALIGN.CENTER,
        tracking=True,
    )


def reverse_label(slide, text, left, top, width):
    add_rect(slide, left, top, width, Inches(0.24), INK)
    add_text(
        slide,
        text,
        left + Inches(0.08),
        top + Inches(0.055),
        width - Inches(0.16),
        Inches(0.12),
        font=SANS,
        size=7.2,
        color=IVORY,
        tracking=True,
    )


def clip_box(slide, left, top, width, height, title, body, *, tag=None, seal=None, title_size=13.5):
    add_rect(slide, left, top, width, height, PANEL, INK, 0.8)
    cursor_y = top + Inches(0.15)
    if tag:
        add_text(
            slide,
            tag,
            left + Inches(0.14),
            cursor_y,
            width - Inches(0.28),
            Inches(0.12),
            font=SANS,
            size=6.8,
            color=CHARCOAL,
            tracking=True,
        )
        cursor_y += Inches(0.22)
    if title:
        add_text(
            slide,
            title,
            left + Inches(0.14),
            cursor_y,
            width - Inches(0.28),
            Inches(0.28),
            font=SERIF,
            size=title_size,
            color=INK,
        )
        cursor_y += Inches(0.38)
    if seal:
        add_oval(slide, left + Inches(0.14), cursor_y, Inches(0.55), None, RED, 1.0)
        add_text(
            slide,
            seal,
            left + Inches(0.17),
            cursor_y + Inches(0.17),
            Inches(0.49),
            Inches(0.18),
            font=SERIF,
            size=9.2,
            color=RED,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            left + Inches(0.82),
            cursor_y + Inches(0.02),
            width - Inches(0.96),
            height - (cursor_y - top) - Inches(0.18),
            font=SANS,
            size=8.6,
            color=DARK_WARM,
            line_spacing=1.22,
        )
        return
    add_text(
        slide,
        body,
        left + Inches(0.14),
        cursor_y,
        width - Inches(0.28),
        height - (cursor_y - top) - Inches(0.14),
        font=SANS,
        size=8.8,
        color=DARK_WARM,
        line_spacing=1.24,
    )


def stat_box(slide, left, top, width, value, label, note):
    add_rect(slide, left, top, width, Inches(0.78), PAPER_LIGHT, RULE, 0.5)
    add_text(slide, value, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.2), font=SERIF, size=17, color=INK)
    add_text(slide, label, left + Inches(0.1), top + Inches(0.34), width - Inches(0.2), Inches(0.12), font=SANS, size=6.4, color=STONE, tracking=True)
    add_text(slide, note, left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), Inches(0.12), font=SANS, size=6.4, color=OLIVE)


def toc_box(slide, left, top, width, item):
    add_rect(slide, left, top, width, Inches(0.78), PAPER_LIGHT, INK, 0.75)
    add_rect(slide, left, top, Inches(0.42), Inches(0.78), INK)
    add_text(slide, item["num"], left, top + Inches(0.18), Inches(0.42), Inches(0.2), font=MONO, size=8, color=IVORY, align=PP_ALIGN.CENTER)
    add_text(slide, item["title"], left + Inches(0.52), top + Inches(0.14), width - Inches(1.1), Inches(0.18), font=SERIF, size=10.5, color=NEAR_BLACK)
    add_text(slide, item["page"], left + width - Inches(0.42), top + Inches(0.18), Inches(0.3), Inches(0.14), font=MONO, size=7.2, color=STONE, align=PP_ALIGN.RIGHT)


def front_page(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_text(slide, spec["masthead"], SHEET_X + Inches(0.36), SHEET_Y + Inches(0.28), Inches(8.7), Inches(0.56), font=SERIF, size=32, color=INK)
    add_text(slide, spec["masthead_sub"], SHEET_X + Inches(0.4), SHEET_Y + Inches(0.78), Inches(8.6), Inches(0.18), font=SANS, size=7.4, color=DARK_WARM, tracking=True)
    draw_issue_mark(slide, SHEET_X + SHEET_W - Inches(1.72), SHEET_Y + Inches(0.26), Inches(1.28), spec["issue_mark_title"], spec["issue_mark_note"])
    add_line(slide, SHEET_X + Inches(0.36), SHEET_Y + Inches(1.08), SHEET_X + SHEET_W - Inches(0.36), SHEET_Y + Inches(1.08), INK, 1.2)
    add_text(slide, spec["dateline_left"], SHEET_X + Inches(0.38), SHEET_Y + Inches(1.18), Inches(5.9), Inches(0.16), font=SANS, size=6.8, color=DARK_WARM, tracking=True)
    add_text(slide, spec["dateline_right"], SHEET_X + SHEET_W - Inches(4.2), SHEET_Y + Inches(1.18), Inches(3.8), Inches(0.16), font=SANS, size=6.8, color=DARK_WARM, tracking=True, align=PP_ALIGN.RIGHT)

    rail_x = SHEET_X + Inches(0.38)
    top = SHEET_Y + Inches(1.56)
    add_rect(slide, rail_x, top, Inches(0.82), Inches(4.84), PANEL, INK, 0.85)
    reverse_label(slide, spec["rail_kicker"], rail_x + Inches(0.06), top + Inches(0.12), Inches(0.7))
    add_text(slide, verticalize(spec["rail_title"]), rail_x + Inches(0.2), top + Inches(0.56), Inches(0.36), Inches(2.75), font=SERIF, size=20, color=INK, align=PP_ALIGN.CENTER, line_spacing=0.92)
    add_text(slide, verticalize(spec["rail_date"]), rail_x + Inches(0.22), top + Inches(3.5), Inches(0.28), Inches(1.4), font=SANS, size=6.8, color=STONE, align=PP_ALIGN.CENTER, line_spacing=1.0)

    lead_x = rail_x + Inches(1.05)
    lead_w = Inches(6.55)
    add_line(slide, lead_x, top + Inches(0.04), lead_x + lead_w, top + Inches(0.04), INK, 0.85)
    reverse_label(slide, spec["lead_label"], lead_x, top + Inches(0.16), Inches(1.42))
    add_text(slide, spec["headline"], lead_x, top + Inches(0.52), lead_w, Inches(0.78), font=SERIF, size=22, color=INK, line_spacing=1.02)
    add_text(slide, spec["lead"], lead_x, top + Inches(1.34), lead_w, Inches(0.52), font=SERIF, size=10.5, color=NEAR_BLACK, line_spacing=1.28)
    add_text(slide, spec["lead_columns"][0], lead_x, top + Inches(1.98), Inches(3.08), Inches(1.28), font=SANS, size=8.8, color=DARK_WARM, line_spacing=1.24)
    add_line(slide, lead_x + Inches(3.21), top + Inches(1.96), lead_x + Inches(3.21), top + Inches(3.3), RULE, 0.5)
    add_text(slide, spec["lead_columns"][1], lead_x + Inches(3.36), top + Inches(1.98), Inches(3.08), Inches(1.28), font=SANS, size=8.8, color=DARK_WARM, line_spacing=1.24)
    add_line(slide, lead_x, top + Inches(3.46), lead_x + lead_w, top + Inches(3.46), RULE, 0.5)
    add_rect(slide, lead_x, top + Inches(3.64), Inches(2.38), Inches(1.1), PAPER_DEEP, INK, 0.7)
    add_line(slide, lead_x - Inches(0.04), top + Inches(4.12), lead_x + Inches(2.18), top + Inches(3.78), RULE_INK, 0.6)
    add_line(slide, lead_x + Inches(1.92), top + Inches(3.58), lead_x + Inches(2.34), top + Inches(4.64), RULE_INK, 0.6)
    add_text(slide, spec["photo_caption"], lead_x + Inches(2.58), top + Inches(3.7), Inches(3.85), Inches(0.92), font=SANS, size=7.2, color=STONE, line_spacing=1.2)

    side_x = lead_x + lead_w + Inches(0.28)
    clip_box(slide, side_x, top, Inches(3.48), Inches(1.28), spec["side_cards"][0]["title"], spec["side_cards"][0]["body"], tag=spec["side_cards"][0]["tag"])
    clip_box(slide, side_x, top + Inches(1.44), Inches(3.48), Inches(1.28), spec["side_cards"][1]["title"], spec["side_cards"][1]["body"], tag=spec["side_cards"][1]["tag"])
    clip_box(slide, side_x, top + Inches(2.88), Inches(3.48), Inches(1.58), "", spec["side_cards"][2]["body"], tag=spec["side_cards"][2]["tag"], seal=spec["side_cards"][2]["seal"])


def component_sheet(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_text(slide, spec["title"], SHEET_X + Inches(0.42), SHEET_Y + Inches(0.42), Inches(6.5), Inches(0.46), font=SERIF, size=24, color=INK)
    add_text(slide, spec["lede"], SHEET_X + Inches(7.1), SHEET_Y + Inches(0.48), Inches(5.0), Inches(0.42), font=SANS, size=8.8, color=DARK_WARM, line_spacing=1.22)
    positions = [0.42, 3.34, 6.26, 9.18]
    for card, offset in zip(spec["cards"], positions, strict=True):
        clip_box(slide, SHEET_X + Inches(offset), SHEET_Y + Inches(1.32), Inches(2.64), Inches(1.48), card["title"], card["body"], tag=card["tag"], title_size=12.2)
    stat_positions = [0.42, 2.9, 5.38]
    for stat, offset in zip(spec["stats"], stat_positions, strict=True):
        stat_box(slide, SHEET_X + Inches(offset), SHEET_Y + Inches(3.1), Inches(2.18), stat["value"], stat["label"], stat["note"])
    timeline_x = SHEET_X + Inches(0.42)
    for idx, item in enumerate(spec["timeline"]):
        x = timeline_x + Inches(2.45) * idx
        add_line(slide, x, SHEET_Y + Inches(4.36), x + Inches(2.05), SHEET_Y + Inches(4.36), INK, 0.8)
        add_text(slide, item["label"], x, SHEET_Y + Inches(4.05), Inches(2.0), Inches(0.18), font=SERIF, size=12.5, color=INK)
        add_text(slide, item["body"], x, SHEET_Y + Inches(4.48), Inches(2.05), Inches(0.42), font=SANS, size=8.0, color=OLIVE, line_spacing=1.22)
    clip_box(
        slide,
        SHEET_X + Inches(8.0),
        SHEET_Y + Inches(3.08),
        Inches(3.52),
        Inches(2.04),
        spec["action_title"],
        spec["action_body"],
        tag="编辑规则",
        title_size=13,
    )


def special_issue(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    reverse_label(slide, spec["tag"], SHEET_X + Inches(0.44), SHEET_Y + Inches(0.42), Inches(2.85))
    add_text(slide, spec["meta"], SHEET_X + SHEET_W - Inches(2.9), SHEET_Y + Inches(0.46), Inches(2.5), Inches(0.22), font=SANS, size=7.2, color=STONE, align=PP_ALIGN.RIGHT, tracking=True)
    add_rect(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(1.02), Inches(7.0), Inches(2.26), INK, INK, 1.0)
    add_rect(slide, SHEET_X + Inches(0.58), SHEET_Y + Inches(1.16), Inches(6.72), Inches(1.98), None, IVORY, 0.7)
    add_text(slide, spec["title"], SHEET_X + Inches(0.78), SHEET_Y + Inches(1.45), Inches(5.95), Inches(0.8), font=SERIF, size=28, color=IVORY, line_spacing=1.02)
    add_text(slide, spec["subtitle"], SHEET_X + Inches(0.82), SHEET_Y + Inches(2.42), Inches(5.9), Inches(0.42), font=SANS, size=8.7, color=IVORY, line_spacing=1.22)
    clip_box(slide, SHEET_X + Inches(7.86), SHEET_Y + Inches(1.02), Inches(2.08), Inches(1.02), "", spec["filed_under"], tag="Filed Under", title_size=11.5)
    clip_box(slide, SHEET_X + Inches(10.08), SHEET_Y + Inches(1.02), Inches(1.48), Inches(1.02), "", spec["version"], tag="Version", title_size=11.5)
    add_line(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(3.56), SHEET_X + Inches(7.42), SHEET_Y + Inches(3.56), INK, 0.8)
    add_text(slide, spec["cover_note"], SHEET_X + Inches(0.44), SHEET_Y + Inches(3.72), Inches(6.9), Inches(0.56), font=SANS, size=8.5, color=DARK_WARM, line_spacing=1.22)
    add_text(slide, spec["cover_meta"], SHEET_X + Inches(8.02), SHEET_Y + Inches(3.72), Inches(3.45), Inches(0.62), font=SANS, size=7.6, color=STONE, line_spacing=1.25, align=PP_ALIGN.RIGHT)
    toc_lefts = [0.44, 3.26, 6.08, 8.9]
    for item, offset in zip(spec["toc"], toc_lefts, strict=True):
        toc_box(slide, SHEET_X + Inches(offset), SHEET_Y + Inches(4.92), Inches(2.52), item)


def article_spread(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_line(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(0.98), SHEET_X + SHEET_W - Inches(0.44), SHEET_Y + Inches(0.98), INK, 0.9)
    add_line(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(1.08), SHEET_X + SHEET_W - Inches(0.44), SHEET_Y + Inches(1.08), RULE_INK, 0.45)
    add_text(slide, spec["chapter"], SHEET_X + Inches(0.46), SHEET_Y + Inches(0.46), Inches(2.2), Inches(0.16), font=SANS, size=7.2, color=STONE, tracking=True)
    add_text(slide, spec["title"], SHEET_X + Inches(0.46), SHEET_Y + Inches(1.2), Inches(6.8), Inches(0.48), font=SERIF, size=23, color=INK)
    add_text(slide, spec["lede"], SHEET_X + Inches(7.3), SHEET_Y + Inches(1.16), Inches(4.35), Inches(0.48), font=SANS, size=8.7, color=DARK_WARM, line_spacing=1.22)
    body_x = SHEET_X + Inches(0.46)
    body_y = SHEET_Y + Inches(1.98)
    add_text(slide, spec["columns"][0], body_x, body_y, Inches(3.35), Inches(2.4), font=SANS, size=9.4, color=CHARCOAL, line_spacing=1.25)
    add_line(slide, body_x + Inches(3.55), body_y, body_x + Inches(3.55), body_y + Inches(2.7), RULE, 0.5)
    add_text(slide, spec["columns"][1], body_x + Inches(3.75), body_y, Inches(3.35), Inches(2.4), font=SANS, size=9.4, color=CHARCOAL, line_spacing=1.25)
    clip_box(
        slide,
        SHEET_X + Inches(8.12),
        SHEET_Y + Inches(1.98),
        Inches(3.48),
        Inches(1.6),
        spec["panel_title"],
        spec["panel_body"],
        tag="SIDE PANEL",
        title_size=12.6,
    )
    for idx, bullet in enumerate(spec["bullets"]):
        add_text(
            slide,
            f"■ {bullet}",
            SHEET_X + Inches(8.22),
            SHEET_Y + Inches(3.88 + idx * 0.28),
            Inches(3.25),
            Inches(0.16),
            font=SANS,
            size=7.8,
            color=OLIVE,
        )
    add_rect(slide, SHEET_X + Inches(8.12), SHEET_Y + Inches(4.84), Inches(3.48), Inches(1.05), PAPER_LIGHT, INK, 0.75)
    add_line(slide, SHEET_X + Inches(8.32), SHEET_Y + Inches(5.0), SHEET_X + Inches(8.32), SHEET_Y + Inches(5.72), INK, 2.4)
    add_text(slide, spec["quote"], SHEET_X + Inches(8.56), SHEET_Y + Inches(4.98), Inches(2.82), Inches(0.42), font=SERIF, size=11.8, color=NEAR_BLACK, line_spacing=1.15)
    add_text(slide, spec["cite"], SHEET_X + Inches(8.56), SHEET_Y + Inches(5.55), Inches(2.82), Inches(0.14), font=SANS, size=6.8, color=STONE, tracking=True)


def correspondence(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_rect(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(0.44), Inches(4.2), Inches(1.02), INK, INK, 1.0)
    add_rect(slide, SHEET_X + Inches(0.58), SHEET_Y + Inches(0.58), Inches(3.92), Inches(0.74), None, IVORY, 0.6)
    add_text(slide, spec["plaque_label"], SHEET_X + Inches(0.74), SHEET_Y + Inches(0.62), Inches(3.3), Inches(0.12), font=SANS, size=7.0, color=IVORY, tracking=True)
    add_text(slide, spec["plaque_value"], SHEET_X + Inches(0.74), SHEET_Y + Inches(0.88), Inches(3.2), Inches(0.28), font=SERIF, size=15.5, color=IVORY)
    clip_box(slide, SHEET_X + Inches(4.9), SHEET_Y + Inches(0.44), Inches(7.44), Inches(1.02), spec["sender_org"], spec["sender_meta"], tag="寄件人", title_size=13.8)

    add_rect(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(1.78), Inches(11.9), Inches(0.98), PANEL, INK, 0.9)
    add_rect(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(1.78), Inches(11.9), Inches(0.26), INK)
    add_text(slide, "函件主题", SHEET_X + Inches(0.58), SHEET_Y + Inches(1.84), Inches(2.2), Inches(0.12), font=SERIF, size=9.4, color=IVORY)
    add_text(slide, "Formal Correspondence", SHEET_X + Inches(8.9), SHEET_Y + Inches(1.84), Inches(3.0), Inches(0.12), font=SANS, size=6.6, color=IVORY, align=PP_ALIGN.RIGHT, tracking=True)
    add_text(slide, spec["subject_title"], SHEET_X + Inches(0.62), SHEET_Y + Inches(2.08), Inches(9.0), Inches(0.24), font=SERIF, size=16.5, color=NEAR_BLACK)
    add_text(slide, f"致：{spec['recipient']}", SHEET_X + Inches(0.62), SHEET_Y + Inches(2.36), Inches(4.6), Inches(0.14), font=SANS, size=7.4, color=STONE)
    add_text(slide, f"分类：{spec['category']}", SHEET_X + Inches(7.05), SHEET_Y + Inches(2.36), Inches(4.8), Inches(0.14), font=SANS, size=7.4, color=STONE, align=PP_ALIGN.RIGHT)

    add_rect(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(2.98), Inches(11.9), Inches(2.95), IVORY, INK, 0.9)
    add_rect(slide, SHEET_X + Inches(0.58), SHEET_Y + Inches(3.12), Inches(11.62), Inches(2.67), None, RULE, 0.45)
    add_text(slide, spec["salutation"], SHEET_X + Inches(0.78), SHEET_Y + Inches(3.18), Inches(3.4), Inches(0.18), font=SERIF, size=11.4, color=NEAR_BLACK)
    para_y = SHEET_Y + Inches(3.48)
    for idx, paragraph in enumerate(spec["paragraphs"]):
        add_text(
            slide,
            paragraph,
            SHEET_X + Inches(0.82),
            para_y + Inches(idx * 0.52),
            Inches(10.95),
            Inches(0.42),
            font=SERIF,
            size=9.5,
            color=DARK_WARM,
            line_spacing=1.22,
        )
    for idx, item in enumerate(spec["evidence"]):
        x = SHEET_X + Inches(0.82 + idx * 3.63)
        add_rect(slide, x, SHEET_Y + Inches(4.98), Inches(3.22), Inches(0.46), PAPER_LIGHT, RULE, 0.45)
        add_text(slide, item, x + Inches(0.12), SHEET_Y + Inches(5.08), Inches(2.98), Inches(0.22), font=SANS, size=7.4, color=OLIVE, line_spacing=1.15)
    add_text(slide, spec["regards"], SHEET_X + Inches(0.82), SHEET_Y + Inches(5.52), Inches(2.2), Inches(0.16), font=SERIF, size=11.2, color=NEAR_BLACK)
    add_text(slide, spec["closing_note"], SHEET_X + Inches(0.82), SHEET_Y + Inches(5.74), Inches(6.2), Inches(0.18), font=SANS, size=7.1, color=STONE)
    add_text(slide, spec["signature"], SHEET_X + Inches(9.02), SHEET_Y + Inches(5.5), Inches(2.8), Inches(0.18), font=SERIF, size=14.5, color=NEAR_BLACK, align=PP_ALIGN.RIGHT)
    add_text(slide, spec["signature_meta"], SHEET_X + Inches(8.66), SHEET_Y + Inches(5.72), Inches(3.2), Inches(0.24), font=SANS, size=7.1, color=STONE, align=PP_ALIGN.RIGHT, line_spacing=1.15)
    add_line(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(6.16), SHEET_X + Inches(12.0), SHEET_Y + Inches(6.16), RULE, 0.45)
    add_text(slide, f"附件：{spec['attachments']}", SHEET_X + Inches(0.44), SHEET_Y + Inches(6.24), Inches(9.8), Inches(0.16), font=SANS, size=7.0, color=STONE)


def routing_desk(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_text(slide, spec["desk_name"], SHEET_X + Inches(0.44), SHEET_Y + Inches(0.44), Inches(5.6), Inches(0.34), font=SERIF, size=24, color=INK)
    add_text(slide, spec["desk_sub"], SHEET_X + Inches(6.4), SHEET_Y + Inches(0.5), Inches(5.5), Inches(0.28), font=SANS, size=8.8, color=DARK_WARM, align=PP_ALIGN.RIGHT)
    top = SHEET_Y + Inches(1.28)
    for idx, wire in enumerate(spec["wires"]):
        y = top + Inches(0.74) * idx
        add_rect(slide, SHEET_X + Inches(0.44), y, Inches(8.12), Inches(0.5), PAPER_LIGHT, INK, 0.65)
        add_text(slide, f"“{wire['prompt']}”", SHEET_X + Inches(0.66), y + Inches(0.14), Inches(5.4), Inches(0.16), font=SERIF, size=11.4, color=NEAR_BLACK)
        add_rect(slide, SHEET_X + Inches(6.18), y + Inches(0.08), Inches(1.18), Inches(0.26), INK)
        add_text(slide, wire["route"], SHEET_X + Inches(6.24), y + Inches(0.13), Inches(1.06), Inches(0.12), font=MONO, size=7.0, color=IVORY, align=PP_ALIGN.CENTER)
        add_text(slide, wire["note"], SHEET_X + Inches(7.56), y + Inches(0.15), Inches(0.82), Inches(0.12), font=SANS, size=6.6, color=STONE, align=PP_ALIGN.RIGHT)
    clip_box(
        slide,
        SHEET_X + Inches(8.92),
        SHEET_Y + Inches(1.28),
        Inches(2.94),
        Inches(2.34),
        "触发原则",
        "\n".join(f"· {item}" for item in spec["routing_rules"]),
        tag="ROUTING RULES",
        title_size=12.6,
    )
    add_rect(slide, SHEET_X + Inches(8.92), SHEET_Y + Inches(4.0), Inches(2.94), Inches(1.6), INK, INK, 0.9)
    draw_issue_mark(slide, SHEET_X + Inches(9.16), SHEET_Y + Inches(4.18), Inches(0.96), spec["stamp_title"], "")
    add_text(slide, spec["stamp_note"], SHEET_X + Inches(10.28), SHEET_Y + Inches(4.34), Inches(1.2), Inches(0.62), font=SANS, size=8.0, color=IVORY, line_spacing=1.18)


def production_desk(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_text(slide, spec["title"], SHEET_X + Inches(0.44), SHEET_Y + Inches(0.44), Inches(6.6), Inches(0.34), font=SERIF, size=24, color=INK)
    add_text(slide, spec["lede"], SHEET_X + Inches(7.12), SHEET_Y + Inches(0.5), Inches(4.8), Inches(0.28), font=SANS, size=8.8, color=DARK_WARM, line_spacing=1.2)
    step_lefts = [0.44, 3.26, 6.08, 8.9]
    for step, offset in zip(spec["steps"], step_lefts, strict=True):
        clip_box(slide, SHEET_X + Inches(offset), SHEET_Y + Inches(1.42), Inches(2.52), Inches(1.5), step["title"], step["body"], tag=step["tag"], title_size=12.4)
    add_rect(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(3.48), Inches(11.42), Inches(0.92), INK, INK, 0.8)
    add_text(slide, "印务命令", SHEET_X + Inches(0.66), SHEET_Y + Inches(3.74), Inches(1.4), Inches(0.12), font=MONO, size=7.0, color=IVORY, tracking=True)
    add_text(slide, "\n".join(spec["commands"]), SHEET_X + Inches(2.02), SHEET_Y + Inches(3.6), Inches(9.46), Inches(0.5), font=MONO, size=8.2, color=IVORY, line_spacing=1.1)
    check_lefts = [0.44, 4.26, 8.08]
    for item, offset in zip(spec["checks"], check_lefts, strict=True):
        add_rect(slide, SHEET_X + Inches(offset), SHEET_Y + Inches(4.78), Inches(3.38), Inches(0.96), PANEL, INK, 0.75)
        add_text(slide, item, SHEET_X + Inches(offset + 0.16), SHEET_Y + Inches(5.1), Inches(3.05), Inches(0.22), font=SERIF, size=11.4, color=NEAR_BLACK, align=PP_ALIGN.CENTER)


def final_edition(prs: Presentation, spec: dict):
    slide = page_base(prs, spec["section"], spec["page"])
    add_text(slide, spec["masthead"], SHEET_X + Inches(0.44), SHEET_Y + Inches(0.44), Inches(3.0), Inches(0.34), font=SERIF, size=26, color=INK)
    add_text(slide, spec["masthead_sub"], SHEET_X + Inches(0.48), SHEET_Y + Inches(0.88), Inches(5.8), Inches(0.14), font=SANS, size=7.2, color=DARK_WARM, tracking=True)
    add_line(slide, SHEET_X + Inches(0.44), SHEET_Y + Inches(1.14), SHEET_X + SHEET_W - Inches(0.44), SHEET_Y + Inches(1.14), INK, 1.1)
    add_text(slide, spec["final_title"], SHEET_X + Inches(0.62), SHEET_Y + Inches(1.66), Inches(6.25), Inches(0.92), font=SERIF, size=28, color=INK, line_spacing=1.02)
    add_text(slide, spec["final_lede"], SHEET_X + Inches(0.66), SHEET_Y + Inches(3.12), Inches(6.2), Inches(1.0), font=SANS, size=9.2, color=CHARCOAL, line_spacing=1.25)
    add_rect(slide, SHEET_X + Inches(8.06), SHEET_Y + Inches(1.52), Inches(3.76), Inches(3.08), INK, INK, 0.9)
    add_text(slide, spec["edition_panel"], SHEET_X + Inches(8.32), SHEET_Y + Inches(1.86), Inches(3.24), Inches(2.42), font=MONO, size=11.2, color=IVORY, line_spacing=1.08, tracking=True)
    add_line(slide, SHEET_X + Inches(0.62), SHEET_Y + Inches(5.02), SHEET_X + Inches(7.1), SHEET_Y + Inches(5.02), RULE_INK, 0.6)
    add_text(slide, spec["closing_line"], SHEET_X + Inches(0.66), SHEET_Y + Inches(5.22), Inches(9.8), Inches(0.3), font=SERIF, size=15.5, color=NEAR_BLACK)


RENDERERS = {
    "front-page": front_page,
    "component-sheet": component_sheet,
    "special-issue": special_issue,
    "article-spread": article_spread,
    "correspondence": correspondence,
    "routing-desk": routing_desk,
    "production-desk": production_desk,
    "final-edition": final_edition,
}


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide in DECK:
        RENDERERS[slide["kind"]](prs, slide)

    prs.save("output.pptx")
    patch_theme_fonts("output.pptx")
    print("✓ Saved output.pptx")


if __name__ == "__main__":
    main()
