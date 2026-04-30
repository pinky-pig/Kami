#!/usr/bin/env python3
"""
Kami republican-manuscript slide deck generator.

Run from this directory via:
  python3 slides.py

The script writes output.pptx. scripts/build.py moves it to
assets/demos/demo-slides.pptx when invoked as `python3 scripts/build.py slides`.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slides_spec import DECK


# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x24, 0x38, 0x51)
NAVY_SOFT = RGBColor(0x4D, 0x5B, 0x6D)
PAPER = RGBColor(0xEB, 0xE5, 0xDD)
PAPER_LIGHT = RGBColor(0xF6, 0xF1, 0xE8)
IVORY = RGBColor(0xF3, 0xEF, 0xEB)
INK = RGBColor(0x23, 0x22, 0x22)
COPY = RGBColor(0x4A, 0x49, 0x47)
MUTED = RGBColor(0x66, 0x63, 0x61)
STONE = RGBColor(0x8B, 0x87, 0x82)
RULE = RGBColor(0xD0, 0xC7, 0xBB)
TAG = RGBColor(0xD5, 0xDE, 0xE7)

SERIF = "display"
SERIF_LATIN = "Newsreader"
SERIF_EA = "京華老宋体"
SERIF_FALLBACK_EA = "Source Han Serif SC"
SANS = "body"
SANS_LATIN = "Newsreader"
SANS_EA = "TsangerJinKai02-W04"
MONO = "mono"
MONO_LATIN = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD = Inches(0.42)
INNER_X = PAD
INNER_Y = PAD
INNER_W = SLIDE_W - PAD * 2
INNER_H = SLIDE_H - PAD * 2
TEXTURE = Path(__file__).resolve().parent.parent / "images" / "paper-overlay.png"
DECK_BY_KIND = {slide["kind"]: slide for slide in DECK}
COVER_PLAQUE_LEFT = INNER_X + Inches(0.72)
COVER_PLAQUE_TOP = INNER_Y + Inches(1.78)
COVER_COPY_LEFT = INNER_X + Inches(8.18)
COVER_COPY_TOP = INNER_Y + Inches(2.02)
COVER_METRIC_TOP = INNER_Y + Inches(3.84)
COVER_META_TOP = INNER_Y + Inches(5.7)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rgb(color: RGBColor) -> RGBColor:
    return color


def asset_image_path(filename: str) -> Path:
    root = Path(__file__).resolve().parent.parent
    candidates = [
        root / "demos" / "mock-assets" / filename,
        root / "images" / filename,
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0]


def visual_slot_copy(slide: dict) -> tuple[str, str, str]:
    diagram = slide.get("diagram")
    label = slide.get("visual_label")
    if not label:
        label = f"DIAGRAM SLOT · {diagram.upper()}" if diagram else "VISUAL SLOT"
    title = slide.get("visual_title") or ("优先使用现有 diagrams" if diagram else "默认纯色占位")
    note = slide.get("visual_note")
    if not note:
        if diagram:
            note = f"优先嵌入 assets/diagrams/{diagram}.html 的 SVG；没有图就继续保留纯色占位。"
        else:
            note = "若需要图表，优先改用 assets/diagrams；若需要真实图片，再填 image。"
    return label, title, note


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, weight=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(weight)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font=SANS,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.0,
    tracking=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text.upper() if tracking else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    apply_typeface(run, font)
    return box


def latin_font(font: str) -> str:
    if font == SERIF:
        return SERIF_LATIN
    if font == SANS:
        return SANS_LATIN
    if font == MONO:
        return MONO_LATIN
    return font


def east_asian_font(font: str) -> str:
    if font == SERIF:
        return SERIF_EA
    if font == SANS:
        return SANS_EA
    if font == MONO:
        return MONO_LATIN
    return SERIF_FALLBACK_EA


def apply_typeface(run, font):
    """Set Latin, East Asian, and complex-script font slots.

    python-pptx only writes <a:latin> for run.font.name. PowerPoint uses
    <a:ea> for Chinese glyphs, so without this the deck silently falls back
    to the default Chinese font even though the XML appears to name a font.
    """
    latin_typeface = latin_font(font)
    # PPTX separates Latin and East Asian font slots. Keep Latin labels in
    # their requested face, and map Chinese glyphs by semantic role so body
    # copy stays readable while display serif keeps the manuscript tone.
    east_asian = east_asian_font(font)
    run.font.name = latin_typeface
    r_pr = run._r.get_or_add_rPr()
    for tag in ("{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
                "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
                "{http://schemas.openxmlformats.org/drawingml/2006/main}cs"):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", latin_typeface)
    ea = OxmlElement("a:ea")
    ea.set("typeface", east_asian)
    cs = OxmlElement("a:cs")
    cs.set("typeface", latin_typeface)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def patch_theme_fonts(pptx_path: str):
    """Make theme defaults match the manuscript font-role split."""
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = text.replace('<a:ea typeface=""/>', f'<a:ea typeface="{SANS_EA}"/>')
                text = text.replace('script="Hans" typeface="宋体"', f'script="Hans" typeface="{SANS_EA}"')
                text = text.replace('script="Hant" typeface="新細明體"', f'script="Hant" typeface="{SANS_EA}"')
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def add_line(slide, left, top, width, color=NAVY, weight=1.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(weight)
    return line


def add_bullets(slide, items, left, top, width, height, size=16, color=COPY):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(7)
        p.line_spacing = 1.18
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        apply_typeface(run, SANS)
    return box


def framed_slide(prs, section="KAMI", page=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, INNER_X, INNER_Y, INNER_W, INNER_H, PAPER_LIGHT)
    if TEXTURE.exists():
        slide.shapes.add_picture(str(TEXTURE), INNER_X, INNER_Y, width=INNER_W, height=INNER_H)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        INNER_X + Inches(0.18),
        INNER_Y + Inches(0.18),
        INNER_W - Inches(0.36),
        INNER_H - Inches(0.36),
        None,
        NAVY,
        1.0,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        INNER_X + Inches(0.29),
        INNER_Y + Inches(0.29),
        INNER_W - Inches(0.58),
        INNER_H - Inches(0.58),
        None,
        NAVY_SOFT,
        0.45,
    )
    add_text(
        slide,
        section,
        INNER_X + Inches(0.48),
        INNER_Y + INNER_H - Inches(0.38),
        Inches(3.4),
        Inches(0.16),
        font=MONO,
        size=7.5,
        color=STONE,
        tracking=True,
    )
    if page is not None:
        add_text(
            slide,
            f"{page:02d} / 08",
            INNER_X + INNER_W - Inches(1.5),
            INNER_Y + INNER_H - Inches(0.38),
            Inches(1.0),
            Inches(0.16),
            font=MONO,
            size=7.5,
            color=STONE,
            align=PP_ALIGN.RIGHT,
        )
    return slide


def title_plaque(slide, left, top, width, height, kicker, title, subtitle=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, NAVY)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.13),
        top + Inches(0.13),
        width - Inches(0.26),
        height - Inches(0.26),
        NAVY,
        IVORY,
        0.75,
    )
    add_text(
        slide,
        kicker,
        left + Inches(0.38),
        top + Inches(0.32),
        width - Inches(0.76),
        Inches(0.28),
        font=MONO,
        size=9,
        color=TAG,
        tracking=True,
    )
    add_text(
        slide,
        title,
        left + Inches(0.36),
        top + Inches(0.76),
        width - Inches(0.72),
        height - Inches(1.15),
        font=SERIF,
        size=34,
        color=IVORY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=0.92,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            left + Inches(0.38),
            top + height - Inches(0.48),
            width - Inches(0.76),
            Inches(0.22),
            font=SANS,
            size=10,
            color=TAG,
        )


def section_header(slide, number, title, lede, page):
    add_text(
        slide,
        f"{number:02d} · SECTION",
        INNER_X + Inches(0.7),
        INNER_Y + Inches(0.58),
        Inches(2.4),
        Inches(0.2),
        font=MONO,
        size=8,
        color=STONE,
        tracking=True,
    )
    add_text(
        slide,
        title,
        INNER_X + Inches(0.7),
        INNER_Y + Inches(0.92),
        Inches(5.7),
        Inches(0.55),
        font=SERIF,
        size=26,
        color=INK,
    )
    add_line(slide, INNER_X + Inches(0.7), INNER_Y + Inches(1.55), Inches(4.2), NAVY, 0.75)
    add_text(
        slide,
        lede,
        INNER_X + Inches(7.0),
        INNER_Y + Inches(0.9),
        Inches(4.8),
        Inches(0.76),
        font=SANS,
        size=11,
        color=MUTED,
        line_spacing=1.14,
    )
    add_text(
        slide,
        f"{page:02d} / 08",
        INNER_X + INNER_W - Inches(1.5),
        INNER_Y + INNER_H - Inches(0.38),
        Inches(1.0),
        Inches(0.16),
        font=MONO,
        size=7.5,
        color=STONE,
        align=PP_ALIGN.RIGHT,
    )


def metric_card(slide, left, top, width, value, label, note):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Inches(1.18), IVORY, NAVY_SOFT, 0.8)
    add_text(slide, value, left + Inches(0.18), top + Inches(0.16), width - Inches(0.36), Inches(0.38), font=SERIF, size=24, color=NAVY)
    add_text(slide, label, left + Inches(0.18), top + Inches(0.58), width - Inches(0.36), Inches(0.18), font=SANS, size=8.5, color=STONE, tracking=True)
    add_text(slide, note, left + Inches(0.18), top + Inches(0.82), width - Inches(0.36), Inches(0.2), font=SANS, size=8.5, color=MUTED)


def archive_card(slide, left, top, width, height, title, body, label=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, IVORY, RULE, 0.7)
    if label:
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.34), NAVY)
        add_text(slide, label, left + Inches(0.18), top + Inches(0.1), width - Inches(0.36), Inches(0.12), font=MONO, size=7.2, color=TAG, tracking=True)
        body_top = top + Inches(0.48)
    else:
        body_top = top + Inches(0.24)
    add_text(slide, title, left + Inches(0.2), body_top, width - Inches(0.4), Inches(0.34), font=SERIF, size=17, color=INK)
    add_text(slide, body, left + Inches(0.2), body_top + Inches(0.48), width - Inches(0.4), height - Inches(0.76), font=SANS, size=10.5, color=COPY, line_spacing=1.15)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_cover(prs):
    slide = DECK_BY_KIND["cover"]
    s = framed_slide(prs, slide["section"], slide["page"])
    title_plaque(
        s,
        COVER_PLAQUE_LEFT,
        COVER_PLAQUE_TOP,
        Inches(7.0),
        Inches(2.8),
        slide["kicker"],
        slide["title"],
        slide["subtitle"],
    )
    add_text(
        s,
        slide["body"],
        COVER_COPY_LEFT,
        COVER_COPY_TOP,
        Inches(3.45),
        Inches(1.46),
        font=SANS,
        size=13,
        color=COPY,
        line_spacing=1.12,
    )
    metric_card(
        s,
        COVER_COPY_LEFT,
        COVER_METRIC_TOP,
        Inches(1.55),
        slide["metrics"][0]["value"],
        slide["metrics"][0]["label"],
        slide["metrics"][0]["note"],
    )
    metric_card(
        s,
        INNER_X + Inches(9.9),
        COVER_METRIC_TOP,
        Inches(1.55),
        slide["metrics"][1]["value"],
        slide["metrics"][1]["label"],
        slide["metrics"][1]["note"],
    )
    add_text(
        s,
        slide["meta"],
        INNER_X + Inches(0.74),
        COVER_META_TOP,
        Inches(4.0),
        Inches(0.22),
        font=MONO,
        size=8.5,
        color=STONE,
        tracking=True,
    )


def slide_principle(prs):
    slide = DECK_BY_KIND["principle"]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    x0 = INNER_X + Inches(0.72)
    y = INNER_Y + Inches(2.05)
    w = Inches(2.52)
    for idx, card in enumerate(slide["cards"]):
        archive_card(s, x0 + (w + Inches(0.35)) * idx, y, w, Inches(2.45), card["title"], card["body"], card["label"])
    add_text(
        s,
        slide["summary"],
        INNER_X + Inches(1.3),
        INNER_Y + Inches(5.26),
        Inches(10.0),
        Inches(0.45),
        font=SERIF,
        size=24,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )


def slide_image_focus(prs, key):
    slide = DECK_BY_KIND[key]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    image_box_left = INNER_X + Inches(0.78)
    image_box_top = INNER_Y + Inches(2.02)
    image_box_width = Inches(6.10)
    image_box_height = Inches(3.34)
    add_shape(s, MSO_SHAPE.RECTANGLE, image_box_left, image_box_top, image_box_width, image_box_height, IVORY, NAVY_SOFT, 0.7)
    image_name = slide.get("image")
    image_path = asset_image_path(image_name) if image_name else None
    if image_path and image_path.exists():
        s.shapes.add_picture(str(image_path), image_box_left, image_box_top, width=image_box_width, height=Inches(2.96))
        add_text(
            s,
            slide["caption"],
            image_box_left + Inches(0.18),
            image_box_top + Inches(3.02),
            image_box_width - Inches(0.36),
            Inches(0.16),
            font=MONO,
            size=7.5,
            color=STONE,
            tracking=True,
        )
    else:
        label, title, note = visual_slot_copy(slide)
        add_shape(s, MSO_SHAPE.RECTANGLE, image_box_left, image_box_top, image_box_width, image_box_height, TAG, NAVY, 0.8)
        add_text(
            s,
            label,
            image_box_left + Inches(0.20),
            image_box_top + Inches(0.20),
            image_box_width - Inches(0.40),
            Inches(0.18),
            font=MONO,
            size=8,
            color=NAVY,
            tracking=True,
        )
        add_text(
            s,
            title,
            image_box_left + Inches(0.28),
            image_box_top + Inches(1.05),
            Inches(4.70),
            Inches(1.05),
            font=SERIF,
            size=27,
            color=INK,
            line_spacing=1.05,
        )
        add_text(
            s,
            note,
            image_box_left + Inches(0.28),
            image_box_top + Inches(2.38),
            image_box_width - Inches(0.60),
            Inches(0.54),
            font=SANS,
            size=11,
            color=COPY,
            line_spacing=1.18,
        )
    metric_card(
        s,
        INNER_X + Inches(7.18),
        INNER_Y + Inches(2.10),
        Inches(2.08),
        slide["metrics"][0]["value"],
        slide["metrics"][0]["label"],
        slide["metrics"][0]["note"],
    )
    metric_card(
        s,
        INNER_X + Inches(9.44),
        INNER_Y + Inches(2.10),
        Inches(2.08),
        slide["metrics"][1]["value"],
        slide["metrics"][1]["label"],
        slide["metrics"][1]["note"],
    )
    add_bullets(
        s,
        slide["bullets"],
        INNER_X + Inches(7.18),
        INNER_Y + Inches(3.72),
        Inches(4.10),
        Inches(1.62),
        size=11,
        color=COPY,
    )


def slide_visual_language(prs):
    slide = DECK_BY_KIND["visual-language"]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    swatch_colors = {
        "#243851": NAVY,
        "#EBE5DD": PAPER,
        "#F3EFEB": IVORY,
        "#D0C7BB": RULE,
    }
    for i, swatch in enumerate(slide["swatches"]):
        left = INNER_X + Inches(0.78) + Inches(2.8) * i
        top = INNER_Y + Inches(2.15)
        color = swatch_colors[swatch["hex"].upper()]
        add_shape(s, MSO_SHAPE.RECTANGLE, left, top, Inches(2.2), Inches(1.15), color, NAVY_SOFT if i else IVORY, 0.5)
        add_text(s, swatch["name"], left, top + Inches(1.35), Inches(2.2), Inches(0.2), font=MONO, size=8, color=STONE, tracking=True, align=PP_ALIGN.CENTER)
        add_text(s, swatch["hex"], left, top + Inches(1.65), Inches(2.2), Inches(0.25), font=SERIF, size=15, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, swatch["use"], left, top + Inches(1.98), Inches(2.2), Inches(0.24), font=SANS, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    archive_card(
        s,
        INNER_X + Inches(1.2),
        INNER_Y + Inches(5.05),
        Inches(10.5),
        Inches(0.88),
        slide["callout"]["title"],
        slide["callout"]["body"],
    )


def slide_templates(prs):
    slide = DECK_BY_KIND["templates"]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    x0 = INNER_X + Inches(0.75)
    y0 = INNER_Y + Inches(2.03)
    for i, card in enumerate(slide["cards"]):
        row = i // 3
        col = i % 3
        archive_card(
            s,
            x0 + Inches(3.8) * col,
            y0 + Inches(1.62) * row,
            Inches(3.36),
            Inches(1.22),
            card["title"],
            card["body"],
            card["label"],
        )


def slide_prompt(prs):
    slide = DECK_BY_KIND["natural-prompts"]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    for i, prompt in enumerate(slide["prompts"]):
        top = INNER_Y + Inches(2.0 + i * 0.72)
        add_shape(s, MSO_SHAPE.RECTANGLE, INNER_X + Inches(1.0), top, Inches(8.7), Inches(0.48), IVORY, RULE, 0.45)
        add_text(s, f"“{prompt['text']}”", INNER_X + Inches(1.25), top + Inches(0.14), Inches(6.6), Inches(0.18), font=SERIF, size=14, color=INK)
        add_text(s, prompt["route"], INNER_X + Inches(8.35), top + Inches(0.15), Inches(1.0), Inches(0.16), font=MONO, size=8, color=NAVY, tracking=True, align=PP_ALIGN.RIGHT)
    add_shape(s, MSO_SHAPE.RECTANGLE, INNER_X + Inches(10.12), INNER_Y + Inches(2.0), Inches(1.4), Inches(3.38), NAVY)
    add_text(s, slide["banner"], INNER_X + Inches(10.27), INNER_Y + Inches(2.98), Inches(1.1), Inches(0.8), font=MONO, size=12, color=IVORY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def slide_density(prs):
    slide = DECK_BY_KIND["density"]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    metric_xs = [Inches(0.95), Inches(3.48), Inches(6.01), Inches(8.54)]
    for x, metric in zip(metric_xs, slide["metrics"], strict=True):
        metric_card(s, INNER_X + x, INNER_Y + Inches(2.08), Inches(2.2), metric["value"], metric["label"], metric["note"])
    add_bullets(
        s,
        slide["bullets"],
        INNER_X + Inches(1.1),
        INNER_Y + Inches(4.25),
        Inches(10.6),
        Inches(1.2),
        size=15,
    )


def slide_pipeline(prs):
    slide = DECK_BY_KIND["delivery"]
    s = framed_slide(prs, slide["section"], slide["page"])
    section_header(
        s,
        slide["number"],
        slide["title"],
        slide["lede"],
        slide["page"],
    )
    y = INNER_Y + Inches(2.7)
    for i, step in enumerate(slide["steps"]):
        x = INNER_X + Inches(0.9) + Inches(2.86) * i
        archive_card(s, x, y, Inches(2.35), Inches(1.28), step["title"], step["body"], step["label"])
        if i < len(slide["steps"]) - 1:
            add_line(s, x + Inches(2.45), y + Inches(0.64), Inches(0.32), NAVY_SOFT, 1.0)
    add_text(
        s,
        "\n".join(slide["commands"]),
        INNER_X + Inches(1.3),
        INNER_Y + Inches(5.2),
        Inches(9.8),
        Inches(0.62),
        font=MONO,
        size=10.5,
        color=NAVY,
        align=PP_ALIGN.CENTER,
        line_spacing=1.05,
    )


def slide_end(prs):
    slide = DECK_BY_KIND["end"]
    s = framed_slide(prs, slide["section"], 8)
    add_text(
        s,
        slide["title"],
        INNER_X + Inches(1.0),
        INNER_Y + Inches(2.4),
        Inches(10.8),
        Inches(0.8),
        font=SERIF,
        size=39,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    add_line(s, INNER_X + Inches(5.6), INNER_Y + Inches(3.5), Inches(1.2), NAVY, 1.3)
    add_text(
        s,
        slide["body"],
        INNER_X + Inches(1.0),
        INNER_Y + Inches(3.95),
        Inches(10.8),
        Inches(0.28),
        font=SANS,
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        slide["meta"],
        INNER_X + Inches(1.0),
        INNER_Y + Inches(5.45),
        Inches(10.8),
        Inches(0.2),
        font=MONO,
        size=8,
        color=STONE,
        tracking=True,
        align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_principle(prs)
    slide_image_focus(prs, "image-focus-factory")
    slide_image_focus(prs, "image-focus-robotaxi")
    slide_image_focus(prs, "image-focus-energy")
    slide_image_focus(prs, "image-focus-optimus")
    slide_templates(prs)
    slide_end(prs)

    prs.save("output.pptx")
    patch_theme_fonts("output.pptx")
    print("✓ Saved output.pptx")


if __name__ == "__main__":
    main()
