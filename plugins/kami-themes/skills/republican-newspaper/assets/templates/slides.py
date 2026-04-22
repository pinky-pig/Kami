#!/usr/bin/env python3
"""
Kami republican-newspaper slide deck generator.

Run from this directory via:
  python3 slides.py

The script writes output.pptx. scripts/build.py moves it to
assets/examples/slides.pptx when invoked as `python3 scripts/build.py slides`.
"""

from __future__ import annotations

import zipfile
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x16, 0x14, 0x11)
NAVY_SOFT = RGBColor(0x4B, 0x40, 0x36)
PAPER = RGBColor(0xE8, 0xDF, 0xC9)
PAPER_LIGHT = RGBColor(0xF8, 0xF1, 0xDD)
IVORY = RGBColor(0xF5, 0xEB, 0xD3)
INK = RGBColor(0x1F, 0x1B, 0x16)
COPY = RGBColor(0x3E, 0x38, 0x2E)
MUTED = RGBColor(0x6A, 0x62, 0x57)
STONE = RGBColor(0x8A, 0x81, 0x73)
RULE = RGBColor(0xB9, 0xAA, 0x91)
TAG = RGBColor(0xE9, 0xDE, 0xC2)

FONT_CONFIG = Path(__file__).resolve().parent.parent / "fonts" / "fonts.json"


def load_font_config() -> dict[str, dict[str, str]]:
    return json.loads(FONT_CONFIG.read_text(encoding="utf-8"))


def font_slot(role: str, slot: str = "latin") -> str:
    return load_font_config()[role][slot]


SERIF = font_slot("serif", "latin")
SERIF_EA = font_slot("serif", "eastAsian")
SANS = font_slot("sans", "latin")
MONO = font_slot("mono", "latin")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD = Inches(0.42)
INNER_X = PAD
INNER_Y = PAD
INNER_W = SLIDE_W - PAD * 2
INNER_H = SLIDE_H - PAD * 2
TEXTURE = Path(__file__).resolve().parent.parent / "images" / "paper-overlay.png"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rgb(color: RGBColor) -> RGBColor:
    return color


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, weight=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(weight)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font=SANS,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.0,
    tracking=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text.upper() if tracking else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    apply_typeface(run, font)
    return box


def apply_typeface(run, font):
    """Set Latin, East Asian, and complex-script font slots.

    python-pptx only writes <a:latin> for run.font.name. PowerPoint uses
    <a:ea> for Chinese glyphs, so without this the deck silently falls back
    to the default Chinese font even though the XML appears to name a font.
    """
    run.font.name = font
    # PPTX separates Latin and East Asian font slots. Keep Latin labels in
    # their requested face, but force Chinese glyphs through the configured
    # East Asian serif face.
    east_asian = SERIF_EA
    r_pr = run._r.get_or_add_rPr()
    for tag in ("{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
                "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
                "{http://schemas.openxmlformats.org/drawingml/2006/main}cs"):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", east_asian)
    cs = OxmlElement("a:cs")
    cs.set("typeface", font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def patch_theme_fonts(pptx_path: str):
    """Make the deck theme default to the configured Chinese serif face."""
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = text.replace('<a:ea typeface=""/>', f'<a:ea typeface="{SERIF_EA}"/>')
                text = text.replace('script="Hans" typeface="宋体"', f'script="Hans" typeface="{SERIF_EA}"')
                text = text.replace('script="Hant" typeface="新細明體"', f'script="Hant" typeface="{SERIF_EA}"')
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def add_line(slide, left, top, width, color=NAVY, weight=1.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(weight)
    return line


def add_bullets(slide, items, left, top, width, height, size=16, color=COPY):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(7)
        p.line_spacing = 1.18
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        apply_typeface(run, SANS)
    return box


def framed_slide(prs, section="KAMI", page=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, INNER_X, INNER_Y, INNER_W, INNER_H, PAPER_LIGHT)
    if TEXTURE.exists():
        slide.shapes.add_picture(str(TEXTURE), INNER_X, INNER_Y, width=INNER_W, height=INNER_H)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        INNER_X + Inches(0.18),
        INNER_Y + Inches(0.18),
        INNER_W - Inches(0.36),
        INNER_H - Inches(0.36),
        None,
        NAVY,
        1.0,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        INNER_X + Inches(0.29),
        INNER_Y + Inches(0.29),
        INNER_W - Inches(0.58),
        INNER_H - Inches(0.58),
        None,
        NAVY_SOFT,
        0.45,
    )
    add_text(
        slide,
        section,
        INNER_X + Inches(0.48),
        INNER_Y + INNER_H - Inches(0.38),
        Inches(3.4),
        Inches(0.16),
        font=MONO,
        size=7.5,
        color=STONE,
        tracking=True,
    )
    if page is not None:
        add_text(
            slide,
            f"{page:02d} / 08",
            INNER_X + INNER_W - Inches(1.5),
            INNER_Y + INNER_H - Inches(0.38),
            Inches(1.0),
            Inches(0.16),
            font=MONO,
            size=7.5,
            color=STONE,
            align=PP_ALIGN.RIGHT,
        )
    return slide


def title_plaque(slide, left, top, width, height, kicker, title, subtitle=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, NAVY)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.13),
        top + Inches(0.13),
        width - Inches(0.26),
        height - Inches(0.26),
        NAVY,
        IVORY,
        0.75,
    )
    add_text(
        slide,
        kicker,
        left + Inches(0.38),
        top + Inches(0.32),
        width - Inches(0.76),
        Inches(0.28),
        font=MONO,
        size=9,
        color=TAG,
        tracking=True,
    )
    add_text(
        slide,
        title,
        left + Inches(0.36),
        top + Inches(0.76),
        width - Inches(0.72),
        height - Inches(1.15),
        font=SERIF,
        size=34,
        color=IVORY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=0.92,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            left + Inches(0.38),
            top + height - Inches(0.48),
            width - Inches(0.76),
            Inches(0.22),
            font=SANS,
            size=10,
            color=TAG,
        )


def section_header(slide, number, title, lede, page):
    add_text(
        slide,
        f"{number:02d} · SECTION",
        INNER_X + Inches(0.7),
        INNER_Y + Inches(0.58),
        Inches(2.4),
        Inches(0.2),
        font=MONO,
        size=8,
        color=STONE,
        tracking=True,
    )
    add_text(
        slide,
        title,
        INNER_X + Inches(0.7),
        INNER_Y + Inches(0.92),
        Inches(5.7),
        Inches(0.55),
        font=SERIF,
        size=26,
        color=INK,
    )
    add_line(slide, INNER_X + Inches(0.7), INNER_Y + Inches(1.55), Inches(4.2), NAVY, 0.75)
    add_text(
        slide,
        lede,
        INNER_X + Inches(7.0),
        INNER_Y + Inches(0.9),
        Inches(4.8),
        Inches(0.76),
        font=SANS,
        size=11,
        color=MUTED,
        line_spacing=1.14,
    )
    add_text(
        slide,
        f"{page:02d} / 08",
        INNER_X + INNER_W - Inches(1.5),
        INNER_Y + INNER_H - Inches(0.38),
        Inches(1.0),
        Inches(0.16),
        font=MONO,
        size=7.5,
        color=STONE,
        align=PP_ALIGN.RIGHT,
    )


def metric_card(slide, left, top, width, value, label, note):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Inches(1.18), IVORY, NAVY_SOFT, 0.8)
    add_text(slide, value, left + Inches(0.18), top + Inches(0.16), width - Inches(0.36), Inches(0.38), font=SERIF, size=24, color=NAVY)
    add_text(slide, label, left + Inches(0.18), top + Inches(0.58), width - Inches(0.36), Inches(0.18), font=SANS, size=8.5, color=STONE, tracking=True)
    add_text(slide, note, left + Inches(0.18), top + Inches(0.82), width - Inches(0.36), Inches(0.2), font=SANS, size=8.5, color=MUTED)


def archive_card(slide, left, top, width, height, title, body, label=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, IVORY, RULE, 0.7)
    if label:
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.34), NAVY)
        add_text(slide, label, left + Inches(0.18), top + Inches(0.1), width - Inches(0.36), Inches(0.12), font=MONO, size=7.2, color=TAG, tracking=True)
        body_top = top + Inches(0.48)
    else:
        body_top = top + Inches(0.24)
    add_text(slide, title, left + Inches(0.2), body_top, width - Inches(0.4), Inches(0.34), font=SERIF, size=17, color=INK)
    add_text(slide, body, left + Inches(0.2), body_top + Inches(0.48), width - Inches(0.4), height - Inches(0.76), font=SANS, size=10.5, color=COPY, line_spacing=1.15)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_cover(prs):
    s = framed_slide(prs, "REPUBLICAN NEWSPAPER EDITION", 1)
    title_plaque(
        s,
        INNER_X + Inches(0.72),
        INNER_Y + Inches(0.78),
        Inches(7.0),
        Inches(2.8),
        "KAMI · 特刊 DEMO",
        "把 AI 文档\n排成报纸号外",
        "黑色铅字 / 旧报纸底 / 剪报框线",
    )
    add_text(
        s,
        "Just tell Claude what you need:\n“帮我生成一份白皮书” / “生成一份项目方案” / “帮我写一份推荐信” / “做一套汇报 slides”",
        INNER_X + Inches(8.18),
        INNER_Y + Inches(1.0),
        Inches(3.45),
        Inches(1.46),
        font=SANS,
        size=13,
        color=COPY,
        line_spacing=1.12,
    )
    metric_card(s, INNER_X + Inches(8.18), INNER_Y + Inches(2.82), Inches(1.55), "05", "demo set", "docs + slides")
    metric_card(s, INNER_X + Inches(9.9), INNER_Y + Inches(2.82), Inches(1.55), "01", "visual rule", "black ink")
    add_text(
        s,
        "2026.04 · Kami Fork",
        INNER_X + Inches(0.74),
        INNER_Y + Inches(5.7),
        Inches(4.0),
        Inches(0.22),
        font=MONO,
        size=8.5,
        color=STONE,
        tracking=True,
    )


def slide_principle(prs):
    s = framed_slide(prs, "METHOD", 2)
    section_header(
        s,
        1,
        "生成原理",
        "好看的本质不是每次重新设计，而是把内容填进稳定模板。",
        2,
    )
    x0 = INNER_X + Inches(0.72)
    y = INNER_Y + Inches(2.05)
    w = Inches(2.52)
    for idx, (title, body) in enumerate(
        [
            ("路由", "先判断语言与文档类型：one-pager、long-doc、letter、slides。"),
            ("整理", "把 raw material 拆成事实、数字、判断和行动，而不是直接堆文字。"),
            ("填充", "使用固定骨架承载内容，避免 AI 每次自由发挥版式。"),
            ("校验", "构建脚本检查页数、字体、占位符与 CSS 约束。"),
        ]
    ):
        archive_card(s, x0 + (w + Inches(0.35)) * idx, y, w, Inches(2.45), title, body, f"STEP {idx + 1:02d}")
    add_text(
        s,
        "固定版式 + design token + 字体层级 + 构建校验",
        INNER_X + Inches(1.3),
        INNER_Y + Inches(5.26),
        Inches(10.0),
        Inches(0.45),
        font=SERIF,
        size=24,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )


def slide_visual_language(prs):
    s = framed_slide(prs, "DESIGN TOKENS", 3)
    section_header(
        s,
        2,
        "民国报纸视觉语言",
        "不做复古海报，不做彩色杂志，只把专业文档整理成可信的报纸特刊。",
        3,
    )
    swatches = [
        ("Black Ink", NAVY, "#161411", "报头 / 框线 / 反白条"),
        ("Newsprint", PAPER, "#E8DFC9", "正文纸面"),
        ("Clipping", IVORY, "#F5EBD3", "剪报块"),
        ("Warm Rule", RULE, "#B9AA91", "细线和分隔"),
    ]
    for i, (name, color, hex_value, use) in enumerate(swatches):
        left = INNER_X + Inches(0.78) + Inches(2.8) * i
        top = INNER_Y + Inches(2.15)
        add_shape(s, MSO_SHAPE.RECTANGLE, left, top, Inches(2.2), Inches(1.15), color, NAVY_SOFT if i else IVORY, 0.5)
        add_text(s, name, left, top + Inches(1.35), Inches(2.2), Inches(0.2), font=MONO, size=8, color=STONE, tracking=True, align=PP_ALIGN.CENTER)
        add_text(s, hex_value, left, top + Inches(1.65), Inches(2.2), Inches(0.25), font=SERIF, size=15, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, use, left, top + Inches(1.98), Inches(2.2), Inches(0.24), font=SANS, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    archive_card(
        s,
        INNER_X + Inches(1.2),
        INNER_Y + Inches(5.05),
        Inches(10.5),
        Inches(0.88),
        "装饰边界",
        "只允许报头、日期线、反白栏目条、剪报框和一枚红印；不加入现代渐变和彩色杂志图。",
    )


def slide_templates(prs):
    s = framed_slide(prs, "TEMPLATES", 4)
    section_header(
        s,
        3,
        "现在能生成什么",
        "中文 v1 的正式能力聚焦三类高频文档，简历和作品集已保留为补充样例。",
        4,
    )
    cards = [
        ("One-Pager", "公司介绍、项目方案、执行摘要。重点是 30 秒抓住核心。", "1 page"),
        ("Long Doc", "白皮书、长文报告、年度总结。强调章节结构与证据链。", "multi"),
        ("Letter", "正式信件、推荐信、推荐函。关系、证据、匹配、推荐。", "1 page"),
        ("Resume", "固定坐标版简历样例，适合验证黑框报纸式布局。", "demo"),
        ("Portfolio", "作品集样例，保留更多视觉叙事与项目卡片。", "demo"),
    ]
    x0 = INNER_X + Inches(0.75)
    y0 = INNER_Y + Inches(2.03)
    for i, (title, body, tag) in enumerate(cards):
        row = i // 3
        col = i % 3
        archive_card(
            s,
            x0 + Inches(3.8) * col,
            y0 + Inches(1.62) * row,
            Inches(3.36),
            Inches(1.22),
            title,
            body,
            tag.upper(),
        )


def slide_prompt(prs):
    s = framed_slide(prs, "NATURAL PROMPTS", 5)
    section_header(
        s,
        4,
        "不用 slash command",
        "用户只要说任务，skill 根据意图选择模板。输出是文档，不是聊天里的格式建议。",
        5,
    )
    prompts = [
        ("帮我生成一份白皮书", "long-doc"),
        ("生成一份项目方案", "one-pager"),
        ("帮我写一份推荐信", "letter"),
        ("帮我把这些内容排版成好看的 PDF", "infer"),
        ("做一套汇报 slides", "slides"),
    ]
    for i, (text, route) in enumerate(prompts):
        top = INNER_Y + Inches(2.0 + i * 0.72)
        add_shape(s, MSO_SHAPE.RECTANGLE, INNER_X + Inches(1.0), top, Inches(8.7), Inches(0.48), IVORY, RULE, 0.45)
        add_text(s, f"“{text}”", INNER_X + Inches(1.25), top + Inches(0.14), Inches(6.6), Inches(0.18), font=SERIF, size=14, color=INK)
        add_text(s, route, INNER_X + Inches(8.35), top + Inches(0.15), Inches(1.0), Inches(0.16), font=MONO, size=8, color=NAVY, tracking=True, align=PP_ALIGN.RIGHT)
    add_shape(s, MSO_SHAPE.RECTANGLE, INNER_X + Inches(10.12), INNER_Y + Inches(2.0), Inches(1.4), Inches(3.38), NAVY)
    add_text(s, "AUTO\nTRIGGER", INNER_X + Inches(10.27), INNER_Y + Inches(2.98), Inches(1.1), Inches(0.8), font=MONO, size=12, color=IVORY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def slide_density(prs):
    s = framed_slide(prs, "LAYOUT RHYTHM", 6)
    section_header(
        s,
        5,
        "排版骨架不变",
        "我们换的是时代气质，不是阅读效率。信息密度、留白与页数约束仍然可控。",
        6,
    )
    metric_card(s, INNER_X + Inches(0.95), INNER_Y + Inches(2.08), Inches(2.2), "1", "one-pager", "严格单页")
    metric_card(s, INNER_X + Inches(3.48), INNER_Y + Inches(2.08), Inches(2.2), "1", "letter", "严格单页")
    metric_card(s, INNER_X + Inches(6.01), INNER_Y + Inches(2.08), Inches(2.2), "3", "demo long-doc", "章节展开")
    metric_card(s, INNER_X + Inches(8.54), INNER_Y + Inches(2.08), Inches(2.2), "500+", "check rules", "字体 / CSS / 占位符")
    add_bullets(
        s,
        [
            "旧纸底不是装饰图，而是稳定色块。",
            "题签和细线负责时代感，正文仍按专业文档阅读。",
            "配置 serif 字体用于报纸气质，功能文字保持清晰。",
        ],
        INNER_X + Inches(1.1),
        INNER_Y + Inches(4.25),
        Inches(10.6),
        Inches(1.2),
        size=15,
    )


def slide_pipeline(prs):
    s = framed_slide(prs, "DELIVERY", 7)
    section_header(
        s,
        6,
        "交付链路",
        "生成不止是写文件，还要让 PDF/PPTX 和预览资产都能被检查。",
        7,
    )
    steps = [
        ("HTML / PPTX", "内容进入固定模板"),
        ("Build", "生成交付文件"),
        ("Verify", "检查页数、字体、占位符"),
        ("Preview", "导出 demo 预览"),
    ]
    y = INNER_Y + Inches(2.7)
    for i, (title, body) in enumerate(steps):
        x = INNER_X + Inches(0.9) + Inches(2.86) * i
        archive_card(s, x, y, Inches(2.35), Inches(1.28), title, body, f"{i + 1:02d}")
        if i < len(steps) - 1:
            add_line(s, x + Inches(2.45), y + Inches(0.64), Inches(0.32), NAVY_SOFT, 1.0)
    add_text(
        s,
        ".venv/bin/python scripts/build.py --verify one-pager\n.venv/bin/python scripts/build.py --check\n.venv/bin/python scripts/build.py slides",
        INNER_X + Inches(1.3),
        INNER_Y + Inches(5.2),
        Inches(9.8),
        Inches(0.62),
        font=MONO,
        size=10.5,
        color=NAVY,
        align=PP_ALIGN.CENTER,
        line_spacing=1.05,
    )


def slide_end(prs):
    s = framed_slide(prs, "END", 8)
    add_text(
        s,
        "好看的本质是稳定",
        INNER_X + Inches(1.0),
        INNER_Y + Inches(2.4),
        Inches(10.8),
        Inches(0.8),
        font=SERIF,
        size=39,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    add_line(s, INNER_X + Inches(5.6), INNER_Y + Inches(3.5), Inches(1.2), NAVY, 1.3)
    add_text(
        s,
        "固定版式 · design token · 字体层级 · 构建校验",
        INNER_X + Inches(1.0),
        INNER_Y + Inches(3.95),
        Inches(10.8),
        Inches(0.28),
        font=SANS,
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "Kami · Republican Manuscript Edition",
        INNER_X + Inches(1.0),
        INNER_Y + Inches(5.45),
        Inches(10.8),
        Inches(0.2),
        font=MONO,
        size=8,
        color=STONE,
        tracking=True,
        align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_principle(prs)
    slide_visual_language(prs)
    slide_templates(prs)
    slide_prompt(prs)
    slide_density(prs)
    slide_pipeline(prs)
    slide_end(prs)

    prs.save("output.pptx")
    patch_theme_fonts("output.pptx")
    print("✓ Saved output.pptx")


if __name__ == "__main__":
    main()
