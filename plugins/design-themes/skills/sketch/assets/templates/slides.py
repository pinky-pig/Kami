#!/usr/bin/env python3
"""Prompt-native Sketch slide deck generator.

This renderer owns the `.pptx` output. The shared slide content lives in
`slides_spec.py`, which is also consumed by the Slidev renderer.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slides_spec import DECK, TOTAL_SLIDES


PAPER = RGBColor(0xFD, 0xFB, 0xF7)
INK = RGBColor(0x2D, 0x2D, 0x2D)
MUTED = RGBColor(0xE5, 0xE0, 0xD8)
ACCENT = RGBColor(0xFF, 0x4D, 0x4D)
BLUE = RGBColor(0x2D, 0x5D, 0xA1)
NOTE = RGBColor(0xFF, 0xF9, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADING_FONT = "Marker Felt"
BODY_FONT = "Noteworthy"
FONT_EA = "LXGW WenKai"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.55)
CARD_FILL = {
    "card": WHITE,
    "note": NOTE,
}
DOT_FILL = {
    "ink": INK,
    "accent": ACCENT,
    "secondary": BLUE,
}


def apply_typeface(run, latin_font=BODY_FONT, ea_font=FONT_EA):
    run.font.name = latin_font
    r_pr = run._r.get_or_add_rPr()
    for tag in (
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
    ):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", latin_font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", ea_font)
    cs = OxmlElement("a:cs")
    cs.set("typeface", latin_font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def patch_theme_fonts(pptx_path: str):
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = text.replace('<a:latin typeface="+mn-lt"/>', f'<a:latin typeface="{BODY_FONT}"/>')
                text = text.replace('<a:latin typeface="+mj-lt"/>', f'<a:latin typeface="{HEADING_FONT}"/>')
                text = text.replace('<a:latin typeface="Calibri"/>', f'<a:latin typeface="{BODY_FONT}"/>')
                text = text.replace('<a:ea typeface=""/>', f'<a:ea typeface="{FONT_EA}"/>')
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def shape(slide, kind, left, top, width, height, fill, line=INK, weight=2.3, rotation=0):
    s = slide.shapes.add_shape(kind, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(weight)
    s.rotation = rotation
    s.shadow.inherit = False
    return s


def text(slide, value, left, top, width, height, *, size=16, color=INK,
         bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
         line_spacing=1.0, latin_font=BODY_FONT, ea_font=FONT_EA):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    apply_typeface(run, latin_font=latin_font, ea_font=ea_font)
    return box


def dots(slide):
    for row in range(8):
        for col in range(16):
            x = Inches(0.35 + col * 0.82)
            y = Inches(0.28 + row * 0.92)
            shape(slide, MSO_SHAPE.OVAL, x, y, Inches(0.045), Inches(0.045), MUTED, line=None)


def note_card(slide, left, top, width, height, *, fill=WHITE, rotation=0, tape=False, tack=False):
    shadow = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.08), top + Inches(0.08), width, height, INK, line=None, rotation=rotation)
    card = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line=INK, weight=2.4, rotation=rotation)
    if tape:
        tape_shape = shape(slide, MSO_SHAPE.RECTANGLE, left + width / 2 - Inches(0.38), top - Inches(0.12), Inches(0.76), Inches(0.18), MUTED, line=INK, weight=1.1, rotation=rotation - 6)
        tape_shape.fill.fore_color.rgb = MUTED
    if tack:
        shape(slide, MSO_SHAPE.OVAL, left + width / 2 - Inches(0.07), top - Inches(0.06), Inches(0.14), Inches(0.14), ACCENT, line=INK, weight=1.4)
    return shadow, card


def scribble(slide, x1, y1, x2, y2, color=BLUE):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(2.4)
    return c


def bg(slide, page, section):
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, PAPER, line=None)
    dots(slide)
    text(slide, f"SKETCH | {section}", M, Inches(0.22), Inches(4.1), Inches(0.2), size=8, color=BLUE, bold=True, latin_font=HEADING_FONT)
    text(slide, f"{page:02d} / {TOTAL_SLIDES:02d}", SLIDE_W - Inches(1.5), Inches(0.22), Inches(0.95), Inches(0.2), size=8, color=INK, bold=True, align=PP_ALIGN.RIGHT, latin_font=HEADING_FONT)


def slide_cover(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, spec["page"], spec["section"])
    note_card(s, M, Inches(1.0), Inches(7.2), Inches(4.55), fill=WHITE, rotation=-2, tape=True)
    text(s, spec["title"], M + Inches(0.42), Inches(1.78), Inches(5.9), Inches(1.2), size=38, bold=True, line_spacing=0.92, latin_font=HEADING_FONT)
    text(s, spec["subtitle"], M + Inches(0.42), Inches(3.48), Inches(5.8), Inches(0.62), size=14, color=INK, line_spacing=1.15)
    note_card(s, Inches(8.8), Inches(1.35), Inches(3.1), Inches(1.55), fill=NOTE, rotation=2, tack=True)
    text(s, spec["sticky"]["label"], Inches(9.2), Inches(1.72), Inches(0.9), Inches(0.25), size=18, color=ACCENT, bold=True, latin_font=HEADING_FONT, align=PP_ALIGN.CENTER)
    text(s, spec["sticky"]["body"], Inches(9.1), Inches(2.05), Inches(2.45), Inches(0.55), size=16, bold=True, latin_font=HEADING_FONT, align=PP_ALIGN.CENTER)
    scribble(s, Inches(7.35), Inches(3.15), Inches(8.95), Inches(2.25))


def slide_tokens(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, spec["page"], spec["section"])
    text(s, spec["title"], M, Inches(0.9), Inches(6.8), Inches(0.4), size=28, bold=True, latin_font=HEADING_FONT)
    for i, card in enumerate(spec["cards"]):
        x = M + Inches(2.95) * i
        note_card(
            s,
            x,
            Inches(2.0),
            Inches(2.35),
            Inches(2.1),
            fill=CARD_FILL[card["fill"]],
            rotation=card["rotation"],
            tape=card["attachment"] == "tape",
            tack=card["attachment"] == "tack",
        )
        text(s, card["label"], x + Inches(0.24), Inches(2.55), Inches(1.9), Inches(0.3), size=16, bold=True, align=PP_ALIGN.CENTER, latin_font=HEADING_FONT)
        shape(s, MSO_SHAPE.OVAL, x + Inches(0.93), Inches(3.23), Inches(0.5), Inches(0.5), DOT_FILL[card["dot"]], line=INK, weight=1.5)
    text(s, spec["summary"], M, Inches(5.35), Inches(8.5), Inches(0.24), size=13, color=INK)


def slide_system(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, spec["page"], spec["section"])
    text(s, spec["title"], M, Inches(0.92), Inches(8.4), Inches(0.42), size=28, bold=True, latin_font=HEADING_FONT)
    for i, card in enumerate(spec["cards"]):
        x = M + Inches(3.9) * i
        note_card(
            s,
            x,
            Inches(1.95),
            Inches(3.0),
            Inches(2.45),
            fill=CARD_FILL[card["fill"]],
            rotation=card["rotation"],
            tape=card["attachment"] == "tape",
            tack=card["attachment"] == "tack",
        )
        text(s, card["label"], x + Inches(0.26), Inches(3.05), Inches(2.45), Inches(0.34), size=18, bold=True, align=PP_ALIGN.CENTER, latin_font=HEADING_FONT)
        text(s, card["body"], x + Inches(0.22), Inches(3.72), Inches(2.48), Inches(0.34), size=11.5, color=INK, align=PP_ALIGN.CENTER)
    scribble(s, Inches(1.85), Inches(5.5), Inches(11.25), Inches(5.25), color=ACCENT)
    text(s, spec["summary"], Inches(2.0), Inches(5.02), Inches(9.4), Inches(0.26), size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, latin_font=HEADING_FONT)


def slide_outputs(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, spec["page"], spec["section"])
    text(s, spec["title"], M, Inches(0.9), Inches(8.8), Inches(0.42), size=28, bold=True, latin_font=HEADING_FONT)
    for i, row in enumerate(spec["rows"]):
        y = Inches(1.9 + i * 1.28)
        note_card(
            s,
            M + Inches(row["offset"]),
            y,
            Inches(10.9),
            Inches(0.82),
            fill=CARD_FILL[row["fill"]],
            rotation=row["rotation"],
            tape=row["attachment"] == "tape",
            tack=row["attachment"] == "tack",
        )
        text(s, row["name"], M + Inches(0.34), y + Inches(0.22), Inches(2.0), Inches(0.2), size=15.5, bold=True, latin_font=HEADING_FONT)
        text(s, row["desc"], M + Inches(2.75), y + Inches(0.23), Inches(6.8), Inches(0.2), size=12.2, color=INK)


def slide_end(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, spec["page"], spec["section"])
    note_card(s, M, Inches(1.25), Inches(7.8), Inches(3.5), fill=WHITE, rotation=-2, tape=True)
    text(s, spec["title"], M + Inches(0.42), Inches(1.95), Inches(6.9), Inches(1.0), size=34, bold=True, line_spacing=0.92, latin_font=HEADING_FONT)
    note_card(s, Inches(9.1), Inches(1.65), Inches(2.45), Inches(1.2), fill=NOTE, rotation=2, tack=True)
    text(s, f"{spec['sticky']['label']}\n{spec['sticky']['body']}", Inches(9.36), Inches(1.96), Inches(1.95), Inches(0.46), size=15, bold=True, align=PP_ALIGN.CENTER, latin_font=HEADING_FONT)
    text(s, spec["summary"], M + Inches(0.46), Inches(3.68), Inches(6.5), Inches(0.26), size=14, color=BLUE, bold=True, latin_font=HEADING_FONT)


RENDERERS = {
    "cover": slide_cover,
    "tokens": slide_tokens,
    "system": slide_system,
    "outputs": slide_outputs,
    "end": slide_end,
}


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for slide in DECK:
        RENDERERS[slide["kind"]](prs, slide)
    out = Path(__file__).resolve().parent / "output.pptx"
    prs.save(out)
    patch_theme_fonts(str(out))
    print("✓ Saved output.pptx")


if __name__ == "__main__":
    main()
