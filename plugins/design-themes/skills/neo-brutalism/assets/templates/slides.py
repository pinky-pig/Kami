#!/usr/bin/env python3
"""Prompt-native Neo-Brutalism slide deck generator.

This renderer owns the `.pptx` output. The shared slide content lives in
`slides_spec.py`, which is also consumed by the Slidev renderer.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slides_spec import DECK, TOTAL_SLIDES


CREAM = RGBColor(0xFF, 0xFD, 0xF5)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
VIOLET = RGBColor(0xC4, 0xB5, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Inter"
FONT_EA = "Source Han Sans SC"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.55)
TONE_FILL = {
    "paper": WHITE,
    "red": RED,
    "yellow": YELLOW,
    "violet": VIOLET,
    "black": BLACK,
}
TEXT_TONE = {
    "ink": BLACK,
    "paper": WHITE,
}


def apply_typeface(run, font: str = FONT):
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    for tag in (
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
    ):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", FONT_EA)
    cs = OxmlElement("a:cs")
    cs.set("typeface", font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def patch_theme_fonts(pptx_path: str):
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text_xml = data.decode("utf-8")
                text_xml = text_xml.replace('<a:latin typeface="+mn-lt"/>', f'<a:latin typeface="{FONT}"/>')
                text_xml = text_xml.replace('<a:latin typeface="+mj-lt"/>', f'<a:latin typeface="{FONT}"/>')
                text_xml = text_xml.replace('<a:latin typeface="Calibri"/>', f'<a:latin typeface="{FONT}"/>')
                text_xml = text_xml.replace('<a:ea typeface=""/>', f'<a:ea typeface="{FONT_EA}"/>')
                text_xml = text_xml.replace('script="Hans" typeface="宋体"', f'script="Hans" typeface="{FONT_EA}"')
                data = text_xml.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def shape(slide, kind, left, top, width, height, fill, line=None, weight=3, rotation=0):
    box = slide.shapes.add_shape(kind, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(weight)
    box.rotation = rotation
    box.shadow.inherit = False
    return box


def text(
    slide,
    value,
    left,
    top,
    width,
    height,
    *,
    size=16,
    color=BLACK,
    bold=True,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    apply_typeface(run)
    return box


def bg(slide, page: int, section: str):
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    for x in [1.0, 2.35, 3.7, 5.05, 6.4, 7.75, 9.1, 10.45, 11.8]:
        shape(slide, MSO_SHAPE.OVAL, Inches(x), Inches(0.25), Inches(0.04), Inches(0.04), BLACK)
        shape(slide, MSO_SHAPE.OVAL, Inches(x), Inches(6.9), Inches(0.04), Inches(0.04), BLACK)
    shape(slide, MSO_SHAPE.RECTANGLE, Inches(-0.5), Inches(0.95), Inches(2.3), Inches(1.15), RED, BLACK, 3, rotation=-8)
    shape(slide, MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(2.3), Inches(5.55), Inches(1.75), Inches(1.0), YELLOW, BLACK, 3, rotation=12)
    text(slide, f"NEO BRUTALISM | {section}", M, Inches(0.25), Inches(4.2), Inches(0.18), size=8)
    text(slide, f"{page:02d} / {TOTAL_SLIDES:02d}", SLIDE_W - Inches(1.6), Inches(0.25), Inches(1.0), Inches(0.18), size=8, align=PP_ALIGN.RIGHT)


def panel(slide, left, top, width, height, *, tone="paper", rotation=0, shadow=True):
    fill = TONE_FILL[tone]
    if shadow:
        shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.12), top + Inches(0.12), width, height, BLACK, None, rotation=rotation)
    return shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill, BLACK, 4, rotation=rotation)


def badge(slide, value, left, top, *, tone="yellow", width=Inches(1.65), rotation=0):
    panel(slide, left, top, width, Inches(0.34), tone=tone, rotation=rotation, shadow=True)
    text(slide, value, left + Inches(0.08), top + Inches(0.09), width - Inches(0.16), Inches(0.12), size=7, align=PP_ALIGN.CENTER)


def slide_cover(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    panel(slide, M, Inches(1.05), Inches(7.6), Inches(3.9), tone=spec["hero_tone"], rotation=-1.2)
    badge(slide, spec["eyebrow"], M + Inches(0.35), Inches(1.38), tone="yellow")
    text(slide, spec["title"], M + Inches(0.45), Inches(1.92), Inches(6.8), Inches(1.55), size=43, line_spacing=0.86)
    text(slide, spec["strapline"], M + Inches(0.5), Inches(3.9), Inches(6.4), Inches(0.25), size=11)
    first, second = spec["stickers"]
    panel(slide, Inches(8.95), Inches(1.35), Inches(3.0), Inches(1.25), tone=first["tone"], rotation=2)
    text(
        slide,
        first["text"],
        Inches(9.28),
        Inches(1.68),
        Inches(2.3),
        Inches(0.55),
        size=20,
        color=TEXT_TONE[first["text_tone"]],
        align=PP_ALIGN.CENTER,
        line_spacing=0.9,
    )
    panel(slide, Inches(8.75), Inches(3.35), Inches(3.25), Inches(1.35), tone=second["tone"], rotation=-3)
    text(
        slide,
        second["text"],
        Inches(9.05),
        Inches(3.68),
        Inches(2.65),
        Inches(0.55),
        size=19,
        color=TEXT_TONE[second["text_tone"]],
        align=PP_ALIGN.CENTER,
        line_spacing=0.9,
    )


def slide_tokens(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    text(slide, spec["title"], M, Inches(1.05), Inches(5.5), Inches(0.5), size=32)
    for index, card in enumerate(spec["cards"]):
        left = M + Inches(2.95) * index
        panel(slide, left, Inches(2.1), Inches(2.42), Inches(2.05), tone=card["tone"], rotation=(-2 + index))
        text(
            slide,
            card["name"],
            left,
            Inches(2.86),
            Inches(2.42),
            Inches(0.28),
            size=19,
            color=TEXT_TONE[card["text_tone"]],
            align=PP_ALIGN.CENTER,
        )
    text(slide, spec["summary"], M, Inches(5.2), Inches(8.6), Inches(0.35), size=15)


def slide_composition(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    text(slide, spec["title"], M, Inches(0.95), Inches(5.5), Inches(0.5), size=32)
    hero = spec["hero"]
    panel(slide, M, Inches(1.8), Inches(5.2), Inches(3.2), tone=hero["tone"], rotation=1)
    text(
        slide,
        hero["text"],
        M + Inches(0.35),
        Inches(2.35),
        Inches(4.4),
        Inches(0.9),
        size=34,
        color=TEXT_TONE[hero["text_tone"]],
        line_spacing=0.82,
    )
    first, second = spec["stickers"]
    panel(slide, Inches(7.1), Inches(1.55), Inches(4.1), Inches(1.25), tone=first["tone"], rotation=-2)
    text(
        slide,
        first["text"],
        Inches(7.38),
        Inches(1.95),
        Inches(3.45),
        Inches(0.25),
        size=18,
        color=TEXT_TONE[first["text_tone"]],
        align=PP_ALIGN.CENTER,
    )
    panel(slide, Inches(7.55), Inches(3.35), Inches(3.7), Inches(1.35), tone=second["tone"], rotation=3)
    text(
        slide,
        second["text"],
        Inches(7.9),
        Inches(3.68),
        Inches(3.0),
        Inches(0.55),
        size=22,
        color=TEXT_TONE[second["text_tone"]],
        align=PP_ALIGN.CENTER,
        line_spacing=0.9,
    )


def slide_outputs(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    text(slide, spec["title"], M, Inches(1.05), Inches(8.8), Inches(0.5), size=30)
    for index, row in enumerate(spec["rows"]):
        top = Inches(2.05 + index * 1.22)
        panel(slide, M, top, Inches(10.6), Inches(0.78), tone=row["tone"], rotation=row["rotation"])
        text(slide, row["name"], M + Inches(0.25), top + Inches(0.22), Inches(2.3), Inches(0.22), size=16)
        text(
            slide,
            row["desc"],
            M + Inches(3.0),
            top + Inches(0.23),
            Inches(6.7),
            Inches(0.22),
            size=13,
            color=TEXT_TONE[row["text_tone"]],
        )


def slide_end(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    panel(slide, M, Inches(1.3), Inches(8.2), Inches(3.0), tone=spec["hero_tone"], rotation=-2)
    text(slide, spec["title"], M + Inches(0.45), Inches(1.85), Inches(7.25), Inches(1.1), size=38, line_spacing=0.86)
    badge_spec = spec["badge"]
    panel(slide, Inches(9.25), Inches(1.55), Inches(2.55), Inches(2.25), tone=badge_spec["tone"], rotation=8)
    text(
        slide,
        badge_spec["text"],
        Inches(9.55),
        Inches(2.05),
        Inches(2.0),
        Inches(0.65),
        size=22,
        color=TEXT_TONE[badge_spec["text_tone"]],
        align=PP_ALIGN.CENTER,
        line_spacing=0.9,
    )


RENDERERS = {
    "cover": slide_cover,
    "tokens": slide_tokens,
    "composition": slide_composition,
    "outputs": slide_outputs,
    "end": slide_end,
}


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for slide in DECK:
        RENDERERS[slide["kind"]](prs, slide)
    out = Path(__file__).resolve().parent / "output.pptx"
    prs.save(out)
    patch_theme_fonts(str(out))


if __name__ == "__main__":
    main()
