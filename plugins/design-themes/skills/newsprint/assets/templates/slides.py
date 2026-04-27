#!/usr/bin/env python3
"""Prompt-native Newsprint slide deck generator.

This renderer owns the `.pptx` output. The shared slide content lives in
`slides_spec.py`, which is also consumed by the Slidev renderer.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slides_spec import DECK, DISPLAY_TOTAL


PAPER = RGBColor(0xF9, 0xF9, 0xF7)
INK = RGBColor(0x11, 0x11, 0x11)
BODY = RGBColor(0x40, 0x40, 0x40)
META = RGBColor(0x73, 0x73, 0x73)
MUTED = RGBColor(0xE5, 0xE5, 0xE0)
RED = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "TsangerJinKai02"
SERIF_EA = "TsangerJinKai02"
SANS = "TsangerJinKai02"
MONO = "TsangerJinKai02"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.45)
TEXTURE = Path(__file__).resolve().parent.parent / "images" / "paper-overlay.png"


def apply_typeface(run, font):
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    for tag in (
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
    ):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", SERIF_EA)
    cs = OxmlElement("a:cs")
    cs.set("typeface", font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def patch_theme_fonts(pptx_path: str):
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = text.replace('<a:latin typeface="+mn-lt"/>', f'<a:latin typeface="{SANS}"/>')
                text = text.replace('<a:latin typeface="+mj-lt"/>', f'<a:latin typeface="{SERIF}"/>')
                text = text.replace('<a:latin typeface="Calibri"/>', f'<a:latin typeface="{SANS}"/>')
                text = text.replace('<a:ea typeface=""/>', f'<a:ea typeface="{SERIF_EA}"/>')
                text = text.replace('<a:cs typeface=""/>', f'<a:cs typeface="{SERIF}"/>')
                text = text.replace('script="Hans" typeface="宋体"', f'script="Hans" typeface="{SERIF_EA}"')
                text = text.replace('script="Hant" typeface="新細明體"', f'script="Hant" typeface="{SERIF_EA}"')
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def rect(slide, left, top, width, height, fill, line=None, weight=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(weight)
    shape.shadow.inherit = False
    return shape


def line(slide, x1, y1, x2, y2, color=INK, weight=1):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(weight)
    return connector


def text(slide, value, left, top, width, height, *, font=SANS, size=14, color=INK,
         bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.0,
         uppercase=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = value.upper() if uppercase else value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    apply_typeface(run, font)
    return box


def base(prs, section, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    if TEXTURE.exists():
        slide.shapes.add_picture(str(TEXTURE), 0, 0, width=SLIDE_W, height=SLIDE_H)
    text(slide, f"VOL. 01 | {section}", M, Inches(0.28), Inches(3.4), Inches(0.18), font=MONO, size=7, color=META, uppercase=True)
    text(slide, "THE NEWSPRINT REVIEW", Inches(4.15), Inches(0.2), Inches(5.0), Inches(0.35), font=SERIF, size=23, align=PP_ALIGN.CENTER)
    text(slide, f"PAGE {page:02d} / {DISPLAY_TOTAL:02d}", SLIDE_W - Inches(2.2), Inches(0.28), Inches(1.75), Inches(0.18), font=MONO, size=7, color=META, align=PP_ALIGN.RIGHT, uppercase=True)
    line(slide, M, Inches(0.72), SLIDE_W - M, Inches(0.72), INK, 4)
    line(slide, M, Inches(6.78), SLIDE_W - M, Inches(6.78), INK, 1)
    return slide


def label(slide, value, left, top, width=Inches(1.8), red=True):
    fill = RED if red else INK
    rect(slide, left, top, width, Inches(0.28), fill)
    text(slide, value, left + Inches(0.08), top + Inches(0.075), width - Inches(0.16), Inches(0.12),
         font=MONO, size=6.5, color=PAPER, uppercase=True)


def column_rule(slide, x, top=Inches(0.9), bottom=Inches(6.62), weight=1):
    line(slide, x, top, x, bottom, INK, weight)


def slide_cover(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    label(s, spec["label"], M, Inches(0.95), Inches(1.05))
    text(s, spec["title"], M, Inches(1.42), Inches(8.15), Inches(2.35),
         font=SERIF, size=55, color=INK, line_spacing=0.86)
    column_rule(s, Inches(8.85), Inches(0.9), Inches(6.62), 1)
    text(s, spec["body"],
         Inches(9.15), Inches(1.18), Inches(3.35), Inches(1.35), font=SERIF, size=15, color=BODY, line_spacing=1.16)
    rect(s, Inches(9.15), Inches(3.0), Inches(3.35), Inches(1.42), INK)
    text(s, spec["edition"], Inches(9.36), Inches(3.24), Inches(2.8), Inches(0.8),
         font=MONO, size=12, color=PAPER, line_spacing=1.05, uppercase=True)
    line(s, M, Inches(4.25), Inches(8.55), Inches(4.25), INK, 1)
    text(s, spec["note"],
         M, Inches(4.52), Inches(8.0), Inches(0.5), font=MONO, size=10, color=META, uppercase=True)


def slide_grid(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    text(s, spec["title"], M, Inches(1.08), Inches(5.4), Inches(0.78), font=SERIF, size=32)
    text(s, spec["body"],
         Inches(6.0), Inches(1.18), Inches(5.95), Inches(0.55), font=SANS, size=12, color=BODY)
    top = Inches(2.12)
    left = M
    cell_w = (SLIDE_W - 2 * M) / 12
    for item in spec["columns"]:
        index = item["index"] - 1
        fill = MUTED if item["tone"] == "muted" else PAPER
        rect(s, left + cell_w * index, top, cell_w, Inches(2.7), fill, INK, 0.8)
        text(s, str(item["index"]), left + cell_w * index + Inches(0.05), top + Inches(2.35), cell_w - Inches(0.1), Inches(0.15), font=MONO, size=7, color=META, align=PP_ALIGN.CENTER)
    label(s, spec["hero_label"], M + Inches(0.15), top + Inches(0.18), Inches(1.35), False)
    label(s, spec["side_label"], M + cell_w * 8 + Inches(0.15), top + Inches(0.18), Inches(1.35), True)
    text(s, spec["note"],
         M, Inches(5.22), Inches(8.3), Inches(0.3), font=MONO, size=10, color=BODY)


def card(slide, left, top, width, height, title, body, tag):
    line(slide, left, top, left + width, top, INK, 1)
    line(slide, left, top, left, top + height, INK, 1)
    line(slide, left + width, top, left + width, top + height, INK, 1)
    line(slide, left, top + height, left + width, top + height, INK, 1)
    text(slide, tag, left + Inches(0.16), top + Inches(0.16), width - Inches(0.32), Inches(0.15), font=MONO, size=6.5, color=RED, uppercase=True)
    text(slide, title, left + Inches(0.16), top + Inches(0.45), width - Inches(0.32), Inches(0.35), font=SERIF, size=17)
    text(slide, body, left + Inches(0.16), top + Inches(0.95), width - Inches(0.32), height - Inches(1.1), font=SANS, size=9.5, color=BODY, line_spacing=1.18)


def slide_components(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    text(s, spec["title"], M, Inches(1.05), Inches(8.7), Inches(0.65), font=SERIF, size=33)
    w = Inches(3.65)
    y = Inches(2.05)
    for idx, item in enumerate(spec["cards"]):
        card(s, M + w * idx, y, w, Inches(1.75), item["title"], item["body"], item["tag"])
    rect(s, M, Inches(4.45), SLIDE_W - 2 * M, Inches(0.46), INK)
    text(s, spec["ticker"],
         M + Inches(0.18), Inches(4.58), SLIDE_W - 2 * M - Inches(0.36), Inches(0.12), font=MONO, size=8.5, color=PAPER, uppercase=True)


def slide_typography(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    label(s, spec["label"], M, Inches(1.02), Inches(1.0))
    text(s, spec["title"], M, Inches(1.42), Inches(6.25), Inches(2.2), font=SERIF, size=47, line_spacing=0.86)
    column_rule(s, Inches(7.05), Inches(1.0), Inches(6.45), 1)
    text(s, spec["body"],
         Inches(7.38), Inches(1.42), Inches(4.9), Inches(1.25), font=SERIF, size=16, color=BODY, line_spacing=1.2)
    text(s, spec["drop_cap"], Inches(7.38), Inches(3.15), Inches(0.65), Inches(0.7), font=SERIF, size=48, color=RED)
    text(s, f" {spec['drop_text']}",
         Inches(8.0), Inches(3.38), Inches(4.3), Inches(0.4), font=SANS, size=11, color=BODY)


def slide_inverted(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    rect(s, M, Inches(1.05), Inches(5.2), Inches(5.35), INK)
    label(s, spec["panel_label"], M + Inches(0.3), Inches(1.35), Inches(1.55), True)
    text(s, spec["title"], M + Inches(0.3), Inches(1.95), Inches(4.4), Inches(1.4), font=SERIF, size=33, color=PAPER, line_spacing=0.9)
    text(s, spec["body"],
         M + Inches(0.3), Inches(4.15), Inches(4.35), Inches(0.8), font=SANS, size=11, color=PAPER, line_spacing=1.18)
    for i, item in enumerate(spec["steps"]):
        x = Inches(6.1) + Inches(2.15) * i
        text(s, item["number"], x, Inches(1.55), Inches(1.4), Inches(0.45), font=SERIF, size=28, color=RED)
        line(s, x, Inches(2.12), x + Inches(1.68), Inches(2.12), INK, 1)
        text(s, item["title"], x, Inches(2.38), Inches(1.7), Inches(0.28), font=SERIF, size=18)
        text(s, item["body"], x, Inches(2.9), Inches(1.7), Inches(0.7), font=SANS, size=9, color=BODY, line_spacing=1.18)


def slide_outputs(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    text(s, spec["title"], M, Inches(1.08), Inches(8.6), Inches(0.7), font=SERIF, size=34)
    y = Inches(2.05)
    for i, row in enumerate(spec["rows"]):
        top = y + Inches(1.08) * i
        line(s, M, top, SLIDE_W - M, top, INK, 1)
        text(s, row["name"], M + Inches(0.15), top + Inches(0.22), Inches(2.0), Inches(0.2), font=MONO, size=9, color=RED, uppercase=True)
        text(s, row["desc"], Inches(2.75), top + Inches(0.18), Inches(8.7), Inches(0.28), font=SANS, size=13, color=BODY)
    line(s, M, y + Inches(3.24), SLIDE_W - M, y + Inches(3.24), INK, 1)


def slide_end(prs, spec):
    s = base(prs, spec["section"], spec["page"])
    text(s, spec["title"], M, Inches(1.25), Inches(8.6), Inches(2.1), font=SERIF, size=50, line_spacing=0.86)
    rect(s, Inches(9.15), Inches(1.3), Inches(3.1), Inches(3.7), INK)
    text(s, spec["panel"],
         Inches(9.43), Inches(1.7), Inches(2.5), Inches(1.2), font=MONO, size=12, color=PAPER, line_spacing=1.1, uppercase=True)
    text(s, spec["tagline"], M, Inches(4.35), Inches(6.4), Inches(0.45), font=SERIF, size=22, color=BODY)


RENDERERS = {
    "cover": slide_cover,
    "grid": slide_grid,
    "components": slide_components,
    "typography": slide_typography,
    "inverted": slide_inverted,
    "outputs": slide_outputs,
    "end": slide_end,
}


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for slide in DECK:
        RENDERERS[slide["kind"]](prs, slide)
    out = Path(__file__).resolve().parent / "output.pptx"
    prs.save(out)
    patch_theme_fonts(str(out))


if __name__ == "__main__":
    main()
