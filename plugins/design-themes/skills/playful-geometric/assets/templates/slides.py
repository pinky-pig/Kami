#!/usr/bin/env python3
"""Prompt-native Playful Geometric slide deck generator.

This renderer owns the `.pptx` output. The shared slide content lives in
`slides_spec.py`, which is also consumed by the Slidev renderer.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slides_spec import DECK, TOTAL_SLIDES


CREAM = RGBColor(0xFF, 0xFD, 0xF5)
INK = RGBColor(0x1E, 0x29, 0x3B)
TEXT = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
PINK = RGBColor(0xF4, 0x72, 0xB6)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
MINT = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Inter"
FONT_EA = "Source Han Sans SC"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.55)

FILL = {
    "white": WHITE,
    "violet": VIOLET,
    "pink": PINK,
    "yellow": YELLOW,
    "mint": MINT,
}
OFFSET = {
    "white": "violet",
    "violet": "pink",
    "pink": "yellow",
    "yellow": "mint",
    "mint": "violet",
}


def apply_typeface(run, font: str = FONT):
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    for tag in (
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
    ):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", FONT_EA)
    cs = OxmlElement("a:cs")
    cs.set("typeface", font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def patch_theme_fonts(pptx_path: str):
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = text.replace('<a:latin typeface="+mn-lt"/>', f'<a:latin typeface="{FONT}"/>')
                text = text.replace('<a:latin typeface="+mj-lt"/>', f'<a:latin typeface="{FONT}"/>')
                text = text.replace('<a:latin typeface="Calibri"/>', f'<a:latin typeface="{FONT}"/>')
                text = text.replace('<a:ea typeface=""/>', f'<a:ea typeface="{FONT_EA}"/>')
                text = text.replace('script="Hans" typeface="宋体"', f'script="Hans" typeface="{FONT_EA}"')
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)


def shape(slide, kind, left, top, width, height, fill, line=INK, weight=2, rotation=0):
    node = slide.shapes.add_shape(kind, left, top, width, height)
    node.fill.solid()
    node.fill.fore_color.rgb = fill
    if line is None:
        node.line.fill.background()
    else:
        node.line.color.rgb = line
        node.line.width = Pt(weight)
    node.rotation = rotation
    node.shadow.inherit = False
    return node


def text(slide, value, left, top, width, height, *, size=16, color=INK,
         bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
         line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    apply_typeface(run)
    return box


def text_color(fill_name: str) -> RGBColor:
    return WHITE if fill_name == "violet" else INK


def bg(slide, page, section):
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, CREAM, line=None)
    shape(slide, MSO_SHAPE.OVAL, Inches(-0.7), Inches(-0.55), Inches(2.3), Inches(2.3), YELLOW, line=None)
    shape(slide, MSO_SHAPE.OVAL, SLIDE_W - Inches(1.55), Inches(0.72), Inches(0.8), Inches(0.8), PINK, line=None)
    shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, SLIDE_W - Inches(1.4), SLIDE_H - Inches(1.2), Inches(0.9), Inches(0.9), MINT, line=None, rotation=18)
    text(slide, f"PLAYFUL GEOMETRIC | {section}", M, Inches(0.28), Inches(4.7), Inches(0.18), size=8, color=MUTED, bold=True)
    text(slide, f"{page:02d} / {TOTAL_SLIDES:02d}", SLIDE_W - Inches(1.6), Inches(0.28), Inches(1.05), Inches(0.18), size=8, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def sticker(slide, left, top, width, height, *, fill_name="white", offset_name: str | None = None, radius=True, rotation=0):
    offset_fill = FILL[offset_name or OFFSET[fill_name]]
    face_fill = FILL[fill_name]
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
          left + Inches(0.12), top + Inches(0.12), width, height, offset_fill, line=INK, weight=2, rotation=rotation)
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                 left, top, width, height, face_fill, line=INK, weight=2.2, rotation=rotation)


def badge(slide, value, left, top, *, fill_name="yellow", width=Inches(1.55)):
    sticker(slide, left, top, width, Inches(0.34), fill_name=fill_name, offset_name="white")
    text(
        slide,
        value,
        left + Inches(0.08),
        top + Inches(0.095),
        width - Inches(0.16),
        Inches(0.12),
        size=7.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_cover(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    sticker(slide, M, Inches(1.05), Inches(7.75), Inches(3.75), fill_name="white", offset_name="violet")
    badge(slide, spec["eyebrow"], M + Inches(0.35), Inches(1.35), fill_name="yellow")
    text(slide, spec["title"], M + Inches(0.45), Inches(1.92), Inches(6.75), Inches(1.55), size=44, color=INK, bold=True, line_spacing=0.9)
    text(slide, spec["subtitle"], M + Inches(0.5), Inches(3.75), Inches(6.3), Inches(0.65), size=13, color=TEXT, line_spacing=1.15)
    first, second = spec["callouts"]
    sticker(slide, Inches(8.95), Inches(1.35), Inches(3.2), Inches(1.45), fill_name=first["fill"])
    text(slide, first["title"], Inches(9.20), Inches(1.70), Inches(2.65), Inches(0.28), size=16, color=text_color(first["fill"]), bold=True, align=PP_ALIGN.CENTER)
    text(slide, first["body"], Inches(9.18), Inches(2.08), Inches(2.72), Inches(0.40), size=11.5, color=text_color(first["fill"]), align=PP_ALIGN.CENTER, line_spacing=1.1)
    sticker(slide, Inches(9.10), Inches(3.10), Inches(2.95), Inches(1.70), fill_name=second["fill"], rotation=-2)
    text(slide, second["title"], Inches(9.40), Inches(3.48), Inches(2.25), Inches(0.56), size=19, color=text_color(second["fill"]), bold=True, align=PP_ALIGN.CENTER, line_spacing=0.95)
    text(slide, second["body"], Inches(9.36), Inches(4.18), Inches(2.36), Inches(0.35), size=11, color=text_color(second["fill"]), align=PP_ALIGN.CENTER)


def slide_tokens(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    text(slide, spec["title"], M, Inches(1.05), Inches(6.5), Inches(0.5), size=30, bold=True)
    for i, card in enumerate(spec["cards"]):
        x = M + Inches(2.95) * i
        sticker(slide, x, Inches(2.05), Inches(2.45), Inches(2.1), fill_name=card["fill"])
        text(
            slide,
            card["label"],
            x,
            Inches(2.85),
            Inches(2.45),
            Inches(0.28),
            size=18,
            color=text_color(card["fill"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    text(slide, spec["summary"], M, Inches(5.15), Inches(8.7), Inches(0.35), size=14, color=TEXT)


def slide_grid(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    text(slide, spec["title"], M, Inches(0.95), Inches(5.8), Inches(0.5), size=30, bold=True)
    left = M
    top = Inches(1.9)
    width = Inches(3.55)
    height = Inches(1.7)
    for i, card in enumerate(spec["cards"]):
        x = left + Inches(3.0) * (i % 2)
        y = top + Inches(2.05) * (i // 2)
        sticker(slide, x, y, width, height, fill_name=card["fill"])
        badge(slide, card["badge"], x + Inches(0.22), y + Inches(0.22), fill_name="violet", width=Inches(0.65))
        text(slide, card["title"], x + Inches(0.28), y + Inches(0.76), width - Inches(0.56), Inches(0.32), size=18, bold=True, align=PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.OVAL, Inches(8.7), Inches(2.0), Inches(2.2), Inches(2.2), YELLOW, line=None)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(4.55), Inches(3.5), Inches(0.75), VIOLET, line=None, rotation=-7)
    text(slide, spec["callout"], Inches(8.05), Inches(3.1), Inches(3.25), Inches(0.45), size=23, bold=True, align=PP_ALIGN.CENTER)


def slide_outputs(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    text(slide, spec["title"], M, Inches(1.05), Inches(8.8), Inches(0.5), size=30, bold=True)
    for i, row in enumerate(spec["rows"]):
        y = Inches(2.0 + i * 1.25)
        sticker(slide, M, y, Inches(10.8), Inches(0.8), fill_name=row["fill"])
        text(slide, row["name"], M + Inches(0.35), y + Inches(0.22), Inches(2.2), Inches(0.22), size=16, bold=True)
        text(slide, row["desc"], M + Inches(3.0), y + Inches(0.23), Inches(6.8), Inches(0.22), size=13, color=TEXT)


def slide_end(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, spec["page"], spec["section"])
    sticker(slide, M, Inches(1.3), Inches(8.2), Inches(3.0), fill_name="white", offset_name="yellow")
    badge(slide, spec["eyebrow"], M + Inches(0.35), Inches(1.58), fill_name="yellow", width=Inches(1.85))
    text(slide, spec["title"], M + Inches(0.45), Inches(2.10), Inches(7.05), Inches(1.1), size=38, bold=True, line_spacing=0.9)
    side = spec["side_note"]
    sticker(slide, Inches(9.15), Inches(1.65), Inches(2.8), Inches(1.55), fill_name=side["fill"], offset_name="pink")
    text(slide, side["title"], Inches(9.42), Inches(2.00), Inches(2.20), Inches(0.24), size=15.5, bold=True, align=PP_ALIGN.CENTER)
    text(slide, side["body"], Inches(9.38), Inches(2.33), Inches(2.28), Inches(0.36), size=11.5, color=TEXT, align=PP_ALIGN.CENTER, line_spacing=1.08)
    shape(slide, MSO_SHAPE.OVAL, Inches(9.35), Inches(3.55), Inches(2.45), Inches(2.45), PINK, line=None)
    shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(9.1), Inches(5.62), Inches(1.1), Inches(1.1), MINT, line=None, rotation=-15)
    text(slide, spec["summary"], M + Inches(0.5), Inches(3.70), Inches(6.8), Inches(0.35), size=16, color=TEXT)


RENDERERS = {
    "cover": slide_cover,
    "tokens": slide_tokens,
    "grid": slide_grid,
    "outputs": slide_outputs,
    "end": slide_end,
}


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for slide in DECK:
        RENDERERS[slide["kind"]](prs, slide)
    out = Path(__file__).resolve().parent / "output.pptx"
    prs.save(out)
    patch_theme_fonts(str(out))
    print("✓ Saved output.pptx")


if __name__ == "__main__":
    main()
