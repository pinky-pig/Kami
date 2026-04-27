#!/usr/bin/env python3
"""Shared Guizang Magazine PPT components."""

from __future__ import annotations

import zipfile
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


PAPER = RGBColor(0xF1, 0xEF, 0xEA)
PAPER_DEEP = RGBColor(0xF6, 0xF3, 0xEE)
IVORY = RGBColor(0xF8, 0xF4, 0xEE)
INK = RGBColor(0x0A, 0x0A, 0x0B)
WARM_DARK = RGBColor(0x18, 0x18, 0x1A)
CHARCOAL = RGBColor(0x40, 0x3B, 0x37)
OLIVE = RGBColor(0x68, 0x63, 0x5F)
STONE = RGBColor(0x8E, 0x87, 0x7F)
RULE = RGBColor(0xD6, 0xD0, 0xC6)
RULE_STRONG = RGBColor(0xC7, 0xBE, 0xB0)
TINT = RGBColor(0xE8, 0xE5, 0xDE)
DARK_CARD = RGBColor(0x14, 0x10, 0x0F)

SERIF = "TsangerJinKai02"
SERIF_EA = "TsangerJinKai02"
SANS = "TsangerJinKai02"
MONO = "TsangerJinKai02"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAGE_X = Inches(0.72)
PAGE_Y = Inches(0.42)
CONTENT_LEFT = PAGE_X
CONTENT_RIGHT = SLIDE_W - PAGE_X
CONTENT_TOP = Inches(0.92)
CONTENT_BOTTOM = SLIDE_H - Inches(0.72)


@dataclass(frozen=True)
class Theme:
    bg: RGBColor
    fg: RGBColor
    muted: RGBColor
    line: RGBColor
    card: RGBColor
    accent: RGBColor


LIGHT = Theme(PAPER, INK, STONE, RULE, IVORY, INK)
DARK = Theme(INK, PAPER, TINT, TINT, DARK_CARD, PAPER)


def theme_for(name: str) -> Theme:
    # Default Guizang output is a single dark mode. Bright pages should only be
    # requested explicitly; the generator must not alternate dark/light pages.
    return LIGHT if name == "light" else DARK


def rgb(color: RGBColor) -> RGBColor:
    return color


def apply_typeface(run, font: str):
    run.font.name = font
    east_asian = SERIF_EA
    r_pr = run._r.get_or_add_rPr()
    for tag in (
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
    ):
        for child in list(r_pr):
            if child.tag == tag:
                r_pr.remove(child)
    latin = OxmlElement("a:latin")
    latin.set("typeface", font)
    ea = OxmlElement("a:ea")
    ea.set("typeface", east_asian)
    cs = OxmlElement("a:cs")
    cs.set("typeface", font)
    r_pr.append(latin)
    r_pr.append(ea)
    r_pr.append(cs)


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, weight=1.0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(weight)
    shape.shadow.inherit = False
    return shape


def add_line(slide, left, top, width, color, weight=0.9):
    line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, left, top, left + width, top)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(weight)
    return line


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    font=SANS,
    size=16,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.0,
    tracking=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text.upper() if tracking else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    apply_typeface(run, font)
    return box


def add_background(slide, theme: Theme):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme.bg)
    add_line(slide, PAGE_X, Inches(0.58), SLIDE_W - PAGE_X * 2, theme.fg if theme is DARK else theme.accent, 1.0)
    add_line(slide, PAGE_X, SLIDE_H - Inches(0.54), SLIDE_W - PAGE_X * 2, theme.line, 0.5)


def add_chrome(
    slide,
    theme: Theme,
    *,
    chrome_left: str,
    chrome_mid: str = "",
    chrome_right: str = "",
    foot_left: str = "",
    foot_right: str = "",
    page_text: str = "",
):
    add_text(
        slide,
        chrome_left,
        PAGE_X,
        Inches(0.18),
        Inches(2.5),
        Inches(0.18),
        font=MONO,
        size=7.4,
        color=theme.muted,
        tracking=True,
    )
    if chrome_mid:
        add_line(slide, PAGE_X + Inches(2.56), Inches(0.29), Inches(0.35), theme.line, 0.6)
        add_text(
            slide,
            chrome_mid,
            PAGE_X + Inches(3.0),
            Inches(0.18),
            Inches(2.4),
            Inches(0.18),
            font=MONO,
            size=7.4,
            color=theme.muted,
            tracking=True,
        )
    if chrome_right:
        add_text(
            slide,
            chrome_right,
            SLIDE_W - PAGE_X - Inches(2.4),
            Inches(0.18),
            Inches(2.4),
            Inches(0.18),
            font=MONO,
            size=7.4,
            color=theme.muted,
            align=PP_ALIGN.RIGHT,
            tracking=True,
        )
    if foot_left:
        add_text(
            slide,
            foot_left,
            PAGE_X,
            SLIDE_H - Inches(0.38),
            Inches(6.8),
            Inches(0.18),
            font=MONO,
            size=7.1,
            color=theme.muted,
            tracking=True,
        )
    if foot_right:
        add_text(
            slide,
            foot_right,
            SLIDE_W - PAGE_X - Inches(3.2),
            SLIDE_H - Inches(0.38),
            Inches(3.2),
            Inches(0.18),
            font=MONO,
            size=7.1,
            color=theme.muted,
            align=PP_ALIGN.RIGHT,
            tracking=True,
        )
    if page_text:
        add_text(
            slide,
            page_text,
            SLIDE_W - PAGE_X - Inches(1.2),
            SLIDE_H - Inches(0.38),
            Inches(1.2),
            Inches(0.18),
            font=MONO,
            size=7.1,
            color=theme.muted,
            align=PP_ALIGN.RIGHT,
            tracking=True,
        )


def new_slide(prs: Presentation, theme_name: str, **chrome_kwargs):
    theme = theme_for(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    add_chrome(slide, theme, **chrome_kwargs)
    return slide, theme


def add_kicker(slide, theme: Theme, text: str, left, top, width):
    add_text(slide, text, left, top, width, Inches(0.22), font=MONO, size=8.2, color=theme.muted, tracking=True)


def add_title(slide, theme: Theme, text: str, left, top, width, *, size=30):
    add_text(slide, text, left, top, width, Inches(0.95), font=SERIF, size=size, color=theme.fg, line_spacing=0.94)


def add_body(slide, theme: Theme, text: str, left, top, width, height, *, size=11.5, color=None):
    add_text(slide, text, left, top, width, height, font=SANS, size=size, color=color or theme.fg, line_spacing=1.16)


def stat_card(slide, theme: Theme, left, top, width, label: str, value: str, note: str):
    add_line(slide, left, top, width, theme.fg, 1.0)
    add_text(slide, label, left, top + Inches(0.12), width, Inches(0.16), font=MONO, size=7.2, color=theme.muted, tracking=True)
    add_text(slide, value, left, top + Inches(0.34), width, Inches(0.42), font=SERIF, size=28, color=theme.fg)
    add_text(slide, note, left, top + Inches(0.82), width, Inches(0.28), font=SANS, size=8.8, color=theme.muted, line_spacing=1.12)


def meta_stack(slide, theme: Theme, left, top, width, label: str, value: str, note: str = ""):
    add_line(slide, left, top, width, theme.line, 0.7)
    add_text(slide, label, left, top + Inches(0.12), width, Inches(0.14), font=MONO, size=7.0, color=theme.muted, tracking=True)
    add_text(slide, value, left, top + Inches(0.3), width, Inches(0.34), font=SERIF, size=16, color=theme.fg)
    if note:
        add_text(slide, note, left, top + Inches(0.68), width, Inches(0.26), font=SANS, size=8.4, color=theme.muted)


def quote_panel(slide, theme: Theme, left, top, width, quote: str, source: str):
    add_line(slide, left, top, width, theme.fg, 1.0)
    add_text(slide, quote, left, top + Inches(0.18), width, Inches(0.9), font=SERIF, size=20, color=theme.fg, line_spacing=0.98)
    add_text(slide, source, left, top + Inches(1.1), width, Inches(0.18), font=MONO, size=7.2, color=theme.muted, tracking=True)


def pillar_card(slide, theme: Theme, left, top, width, index: str, title: str, body: str):
    add_line(slide, left, top, width, theme.fg, 1.0)
    add_text(slide, index, left, top + Inches(0.12), width, Inches(0.14), font=MONO, size=7.1, color=theme.muted, tracking=True)
    add_text(slide, title, left, top + Inches(0.32), width, Inches(0.42), font=SERIF, size=17, color=theme.fg, line_spacing=0.96)
    add_text(slide, body, left, top + Inches(0.82), width, Inches(0.58), font=SANS, size=8.9, color=theme.muted, line_spacing=1.12)


def pipeline_step(slide, theme: Theme, left, top, width, index: str, title: str, body: str):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top + Inches(0.02), Inches(0.16), Inches(0.16), theme.fg)
    add_text(slide, index, left + Inches(0.28), top, Inches(0.5), Inches(0.16), font=MONO, size=7.1, color=theme.muted, tracking=True)
    add_text(slide, title, left, top + Inches(0.28), width, Inches(0.28), font=SERIF, size=15, color=theme.fg)
    add_text(slide, body, left, top + Inches(0.6), width, Inches(0.44), font=SANS, size=8.8, color=theme.muted, line_spacing=1.12)


def rowline(slide, theme: Theme, left, top, width, title: str, body: str, tag: str):
    add_line(slide, left, top, width, theme.line, 0.7)
    add_text(slide, title, left, top + Inches(0.12), Inches(2.2), Inches(0.2), font=SERIF, size=15, color=theme.fg)
    add_text(slide, body, left + Inches(2.4), top + Inches(0.11), width - Inches(3.8), Inches(0.25), font=SANS, size=8.8, color=theme.muted, line_spacing=1.1)
    add_text(slide, tag, left + width - Inches(1.1), top + Inches(0.12), Inches(1.1), Inches(0.18), font=MONO, size=6.9, color=theme.muted, align=PP_ALIGN.RIGHT, tracking=True)


def image_panel(slide, theme: Theme, left, top, width, height, label: str, caption: str):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, theme.card, theme.line, 0.7)
    add_text(slide, label, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.26), font=MONO, size=8.0, color=theme.muted, align=PP_ALIGN.CENTER, tracking=True)
    add_text(slide, "+", left, top + height / 2 - Inches(0.28), width, Inches(0.3), font=SERIF, size=22, color=theme.muted, align=PP_ALIGN.CENTER)
    add_text(slide, caption, left, top + height + Inches(0.08), width, Inches(0.18), font=MONO, size=6.9, color=theme.muted, align=PP_ALIGN.CENTER, tracking=True)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def patch_theme_fonts(pptx_path: str):
    path = Path(pptx_path)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                text = re.sub(
                    r'<a:(latin|ea|cs) typeface="[^"]*"',
                    lambda m: f'<a:{m.group(1)} typeface="{SERIF_EA}"',
                    text,
                )
                text = re.sub(
                    r'(<a:font script="(?:Hans|Hant|Jpan|Hang|Kore)" typeface=")[^"]*(")',
                    rf'\1{SERIF_EA}\2',
                    text,
                )
                data = text.encode("utf-8")
            dst.writestr(item, data)
    tmp.replace(path)
